#!/usr/bin/env python3
"""
KLA Thermal Calibration Presentation Generator - Fully Editable Version
Creates a modern, visually stunning presentation where all elements are editable
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Circle
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import plotly.io as pio
import os

# Set modern style
plt.style.use('seaborn-v0_8')
sns.set_palette("husl")
pio.renderers.default = "png"

# KLA Brand Colors
KLA_BLUE = RGBColor(0, 51, 102)
KLA_ORANGE = RGBColor(255, 102, 0)
KLA_GRAY = RGBColor(128, 128, 128)
KLA_LIGHT_BLUE = RGBColor(0, 102, 204)

def create_editable_title_slide(prs):
    """Create title slide with editable elements"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape (editable)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = KLA_BLUE
    background.line.fill.background()
    
    # Title textbox (fully editable)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Automated Thermal Calibration Module\nfor Thin-Film Deposition"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle textbox (fully editable)
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Varad Lad | Manufacturing Design Engineer Candidate"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # KLA SensArray Division textbox (fully editable)
    kla_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    kla_frame = kla_box.text_frame
    kla_frame.text = "KLA Corporation | SensArray Division"
    kla_para = kla_frame.paragraphs[0]
    kla_para.font.size = Pt(20)
    kla_para.font.color.rgb = KLA_ORANGE
    kla_para.alignment = PP_ALIGN.CENTER

def create_editable_introduction_slide(prs):
    """Create introduction slide with editable content"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title (editable)
    title = slide.shapes.title
    title.text = "Engineering Excellence in Thermal Systems"
    title.text_frame.paragraphs[0].font.color.rgb = KLA_BLUE
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content placeholder (editable)
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.clear()
    
    # Add editable expertise points
    expertise_points = [
        "5+ years experience in thin-film/vacuum deposition equipment",
        "Expertise in CVD, PVD, ALD techniques and process optimization",
        "Advanced thermal management and sensor-based feedback systems",
        "Proven track record: 28% yield improvement, 12% film uniformity gain",
        "Strong background in semiconductor manufacturing and quality control"
    ]
    
    for point in expertise_points:
        p = content_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = KLA_GRAY
        p.space_after = Pt(12)

def create_editable_problem_slide(prs):
    """Create problem statement slide with editable elements"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "The Challenge: Thin-Film Uniformity in Semiconductor Manufacturing"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable comparison boxes
    # Before box
    before_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(2)
    )
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RGBColor(255, 200, 200)  # Light red
    before_box.line.color.rgb = RGBColor(255, 0, 0)
    before_box.line.width = Pt(3)
    
    before_text = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(2))
    before_text.text_frame.text = "BEFORE:\nNon-Uniform Deposition\n• Temperature variations\n• Inconsistent film thickness\n• Yield losses"
    before_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    before_text.text_frame.paragraphs[0].font.bold = True
    before_text.text_frame.paragraphs[0].font.size = Pt(14)
    
    # After box
    after_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2), Inches(4), Inches(2)
    )
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = RGBColor(200, 255, 200)  # Light green
    after_box.line.color.rgb = RGBColor(0, 255, 0)
    after_box.line.width = Pt(3)
    
    after_text = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(2))
    after_text.text_frame.text = "AFTER:\nOptimized Uniformity\n• Controlled temperature\n• Consistent film thickness\n• Improved yield"
    after_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    after_text.text_frame.paragraphs[0].font.bold = True
    after_text.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Add key challenges textbox (editable)
    challenges_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    challenges_frame = challenges_box.text_frame
    challenges_frame.text = "Key Challenges:\n• Temperature variations causing non-uniform deposition\n• Manual calibration processes leading to inconsistencies\n• Lack of real-time thermal feedback and control\n• Yield losses due to thermal drift during processing"
    challenges_para = challenges_frame.paragraphs[0]
    challenges_para.font.size = Pt(16)
    challenges_para.font.color.rgb = KLA_GRAY

def create_editable_design_slide(prs):
    """Create design approach slide with editable system diagram"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Design Approach: Automated Thermal Calibration System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable shapes for system diagram
    # Thermal Sensors
    sensor_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), Inches(2.5), Inches(2), Inches(1)
    )
    sensor_shape.fill.solid()
    sensor_shape.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
    sensor_text = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(0.5))
    sensor_text.text_frame.text = "Thermal\nSensors"
    sensor_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sensor_text.text_frame.paragraphs[0].font.bold = True
    
    # Control Unit
    control_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.5), Inches(2), Inches(1)
    )
    control_shape.fill.solid()
    control_shape.fill.fore_color.rgb = RGBColor(144, 238, 144)  # Light green
    control_text = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(2), Inches(0.5))
    control_text.text_frame.text = "Control\nUnit"
    control_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    control_text.text_frame.paragraphs[0].font.bold = True
    
    # Heating Elements
    heating_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), Inches(4.5), Inches(2), Inches(1)
    )
    heating_shape.fill.solid()
    heating_shape.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
    heating_text = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(2), Inches(0.5))
    heating_text.text_frame.text = "Heating\nElements"
    heating_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    heating_text.text_frame.paragraphs[0].font.bold = True
    
    # Deposition Chamber
    chamber_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(2), Inches(1)
    )
    chamber_shape.fill.solid()
    chamber_shape.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Orange
    chamber_text = slide.shapes.add_textbox(Inches(4), Inches(5), Inches(2), Inches(0.5))
    chamber_text.text_frame.text = "Deposition\nChamber"
    chamber_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    chamber_text.text_frame.paragraphs[0].font.bold = True
    
    # Add arrows (editable lines)
    # Sensors to Control
    arrow1 = slide.shapes.add_connector(1, Inches(3), Inches(3), Inches(4), Inches(3))
    arrow1.line.color.rgb = RGBColor(255, 0, 0)
    arrow1.line.width = Pt(3)
    
    # Control to Heating
    arrow2 = slide.shapes.add_connector(1, Inches(4), Inches(3.5), Inches(2), Inches(4.5))
    arrow2.line.color.rgb = RGBColor(255, 0, 0)
    arrow2.line.width = Pt(3)
    
    # Heating to Chamber
    arrow3 = slide.shapes.add_connector(1, Inches(3), Inches(5), Inches(4), Inches(5))
    arrow3.line.color.rgb = RGBColor(255, 0, 0)
    arrow3.line.width = Pt(3)

def create_editable_technical_slide(prs):
    """Create technical implementation slide with editable process flow"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Technical Implementation: CVD/PVD/ALD Process Integration"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable process flow boxes
    steps = ['Substrate\nPreparation', 'Thermal\nCalibration', 'Deposition\nProcess', 'Real-time\nMonitoring', 'Quality\nControl']
    x_positions = [1, 3, 5, 7, 9]
    
    for i, (step, x) in enumerate(zip(steps, x_positions)):
        # Add rounded rectangle
        box_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x-0.8), Inches(2), Inches(1.6), Inches(1)
        )
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
        box_shape.line.color.rgb = RGBColor(0, 0, 255)  # Blue border
        
        # Add text
        text_box = slide.shapes.add_textbox(Inches(x-0.8), Inches(2), Inches(1.6), Inches(1))
        text_box.text_frame.text = step
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_box.text_frame.paragraphs[0].font.bold = True
        text_box.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Add arrow to next step
        if i < len(x_positions) - 1:
            arrow = slide.shapes.add_connector(1, Inches(x+0.8), Inches(2.5), Inches(x_positions[i+1]-0.8), Inches(2.5))
            arrow.line.color.rgb = RGBColor(255, 0, 0)
            arrow.line.width = Pt(2)
    
    # Add technical details textbox (editable)
    details_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
    details_frame = details_box.text_frame
    details_frame.text = "Key Technical Features:\n• Multi-zone thermal control with PID feedback loops\n• Real-time temperature monitoring using RTD sensors\n• Automated calibration algorithms for process optimization\n• Integration with existing CVD/PVD/ALD equipment\n• Advanced data logging and analysis capabilities"
    details_para = details_frame.paragraphs[0]
    details_para.font.size = Pt(16)
    details_para.font.color.rgb = KLA_GRAY

def create_editable_analysis_slide(prs):
    """Create analysis slide with editable content"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Analysis & Optimization: DOE Methodology & SPC Control"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable analysis boxes
    # DOE Box
    doe_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(2)
    )
    doe_box.fill.solid()
    doe_box.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
    doe_text = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(2))
    doe_text.text_frame.text = "Design of Experiments (DOE):\n• Temperature optimization\n• Pressure control\n• Gas flow rates\n• Process parameters"
    doe_text.text_frame.paragraphs[0].font.bold = True
    doe_text.text_frame.paragraphs[0].font.size = Pt(14)
    
    # SPC Box
    spc_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5), Inches(2), Inches(4), Inches(2)
    )
    spc_box.fill.solid()
    spc_box.fill.fore_color.rgb = RGBColor(144, 238, 144)  # Light green
    spc_text = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(2))
    spc_text.text_frame.text = "Statistical Process Control (SPC):\n• Real-time monitoring\n• Control charts\n• Process capability\n• Quality metrics"
    spc_text.text_frame.paragraphs[0].font.bold = True
    spc_text.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Add optimization results textbox (editable)
    results_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    results_frame = results_box.text_frame
    results_frame.text = "Optimization Results:\n• Temperature optimization improved uniformity by 12%\n• Pressure control enhanced deposition rate by 15%\n• SPC implementation reduced process variability by 25%\n• Automated calibration achieved 95% consistency vs 72% manual"
    results_para = results_frame.paragraphs[0]
    results_para.font.size = Pt(16)
    results_para.font.color.rgb = KLA_GRAY

def create_editable_results_slide(prs):
    """Create results slide with editable metrics"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Results & Impact: Quantifiable Performance Improvements"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable performance metrics boxes
    metrics = [
        ("Yield Improvement", "28%", RGBColor(255, 165, 0)),  # Orange
        ("Film Uniformity", "12%", RGBColor(0, 255, 0)),      # Green
        ("Process Efficiency", "23%", RGBColor(0, 102, 204)),  # Blue
        ("Cost Savings", "25%", RGBColor(255, 102, 0))        # KLA Orange
    ]
    
    for i, (metric, value, color) in enumerate(metrics):
        x_pos = 0.5 + i * 2.25
        # Metric box
        metric_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(2), Inches(2), Inches(1.5)
        )
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = color
        metric_box.line.color.rgb = RGBColor(0, 0, 0)
        metric_box.line.width = Pt(2)
        
        # Metric text
        metric_text = slide.shapes.add_textbox(Inches(x_pos), Inches(2), Inches(2), Inches(1.5))
        metric_text.text_frame.text = f"{metric}\n{value}"
        metric_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        metric_text.text_frame.paragraphs[0].font.bold = True
        metric_text.text_frame.paragraphs[0].font.size = Pt(14)
        metric_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add key achievements textbox (editable)
    achievements_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.5))
    achievements_frame = achievements_box.text_frame
    achievements_frame.text = "Key Achievements:\n• 28% yield improvement through automated thermal control\n• 12% increase in film uniformity using real-time feedback\n• 23% improvement in process efficiency with optimized parameters\n• 25% cost savings through reduced rework and improved quality\n• 90% defect detection rate with advanced metrology"
    achievements_para = achievements_frame.paragraphs[0]
    achievements_para.font.size = Pt(16)
    achievements_para.font.color.rgb = KLA_GRAY

def create_editable_quality_slide(prs):
    """Create quality assurance slide with editable content"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Quality Assurance: Advanced Metrology & Defect Detection"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable metrology boxes
    metrology_methods = [
        ("SEM Analysis", "Surface morphology and defect detection", RGBColor(173, 216, 230)),
        ("XRD Pattern", "Crystalline structure analysis", RGBColor(144, 238, 144)),
        ("UV-Vis Spectroscopy", "Optical properties measurement", RGBColor(255, 165, 0)),
        ("Defect Detection", "90% detection rate achieved", RGBColor(255, 102, 0))
    ]
    
    for i, (method, description, color) in enumerate(metrology_methods):
        y_pos = 2 + i * 1.2
        # Method box
        method_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(4), Inches(1)
        )
        method_box.fill.solid()
        method_box.fill.fore_color.rgb = color
        method_box.line.color.rgb = RGBColor(0, 0, 0)
        method_box.line.width = Pt(2)
        
        # Method text
        method_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(4), Inches(1))
        method_text.text_frame.text = f"{method}\n{description}"
        method_text.text_frame.paragraphs[0].font.bold = True
        method_text.text_frame.paragraphs[0].font.size = Pt(12)
        method_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Add quality metrics textbox (editable)
    quality_box = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(4))
    quality_frame = quality_box.text_frame
    quality_frame.text = "Quality Metrics:\n• 90% defect detection rate\n• 99.97% film purity achieved\n• <1% thickness variation\n• Real-time quality monitoring\n• Automated defect classification"
    quality_para = quality_frame.paragraphs[0]
    quality_para.font.size = Pt(16)
    quality_para.font.color.rgb = KLA_GRAY

def create_editable_kla_alignment_slide(prs):
    """Create KLA alignment slide with editable content"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "KLA SensArray Alignment: In-Situ Thermal Process Control"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable alignment diagram
    # KLA Mission box
    kla_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(1)
    )
    kla_box.fill.solid()
    kla_box.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Orange
    kla_text = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    kla_text.text_frame.text = "KLA SensArray Mission:\nIn-Situ Thermal Process Control"
    kla_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    kla_text.text_frame.paragraphs[0].font.bold = True
    kla_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # My Experience box
    exp_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3), Inches(1)
    )
    exp_box.fill.solid()
    exp_box.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue
    exp_text = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3), Inches(1))
    exp_text.text_frame.text = "My Experience:\nThermal Calibration & Control"
    exp_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    exp_text.text_frame.paragraphs[0].font.bold = True
    exp_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add alignment points (editable textboxes)
    alignment_points = [
        "Real-time Thermal Monitoring",
        "Process Optimization", 
        "Yield Improvement",
        "Quality Control",
        "APC Integration"
    ]
    
    for i, point in enumerate(alignment_points):
        y_pos = 3.5 + i * 0.6
        # Left side
        left_text = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(3), Inches(0.5))
        left_text.text_frame.text = point
        left_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        left_text.text_frame.paragraphs[0].font.bold = True
        
        # Right side
        right_text = slide.shapes.add_textbox(Inches(6), Inches(y_pos), Inches(3), Inches(0.5))
        right_text.text_frame.text = "✓"
        right_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        right_text.text_frame.paragraphs[0].font.bold = True
        right_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 0)
        
        # Connection line
        line = slide.shapes.add_connector(1, Inches(4), Inches(y_pos+0.25), Inches(6), Inches(y_pos+0.25))
        line.line.color.rgb = RGBColor(128, 128, 128)
        line.line.width = Pt(2)

def create_editable_future_slide(prs):
    """Create future applications slide with editable content"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Future Applications: Scalability to KLA Manufacturing Environment"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = KLA_BLUE
    
    # Add editable scalability diagram
    # Current Implementation
    current_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(1)
    )
    current_box.fill.solid()
    current_box.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
    current_text = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    current_text.text_frame.text = "Current Implementation"
    current_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    current_text.text_frame.paragraphs[0].font.bold = True
    
    # KLA Environment
    kla_env_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(3), Inches(1)
    )
    kla_env_box.fill.solid()
    kla_env_box.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Orange
    kla_env_text = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3), Inches(1))
    kla_env_text.text_frame.text = "KLA Manufacturing Environment"
    kla_env_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    kla_env_text.text_frame.paragraphs[0].font.bold = True
    
    # Add scalability benefits textbox (editable)
    benefits_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(3))
    benefits_frame = benefits_box.text_frame
    benefits_frame.text = "Scalability Benefits:\n• Modular design allows easy integration with existing KLA equipment\n• Real-time thermal control enhances SensArray's in-situ monitoring capabilities\n• Automated calibration reduces manual intervention and improves consistency\n• Advanced analytics support APC (Automated Process Control) implementation\n• Proven methodology can be adapted to various thin-film processes"
    benefits_para = benefits_frame.paragraphs[0]
    benefits_para.font.size = Pt(16)
    benefits_para.font.color.rgb = KLA_GRAY

def create_editable_qa_slide(prs):
    """Create Q&A slide with editable elements"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape (editable)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = KLA_BLUE
    background.line.fill.background()
    
    # Title textbox (editable)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!\nQuestions & Discussion"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Contact info textbox (editable)
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Varad Lad | vlad3@asu.edu | 602-388-6861"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(24)
    contact_para.font.color.rgb = RGBColor(255, 255, 255)
    contact_para.alignment = PP_ALIGN.CENTER

def main():
    """Create the complete editable presentation"""
    print("Creating KLA Thermal Calibration Presentation - Fully Editable Version...")
    
    # Create presentation
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create slides
    create_editable_title_slide(prs)
    create_editable_introduction_slide(prs)
    create_editable_problem_slide(prs)
    create_editable_design_slide(prs)
    create_editable_technical_slide(prs)
    create_editable_analysis_slide(prs)
    create_editable_results_slide(prs)
    create_editable_quality_slide(prs)
    create_editable_kla_alignment_slide(prs)
    create_editable_future_slide(prs)
    create_editable_qa_slide(prs)
    
    # Save presentation
    prs.save('Varad_Lad_KLA_Editable_Presentation.pptx')
    
    print("Editable presentation created successfully: Varad_Lad_KLA_Editable_Presentation.pptx")
    print(f"Total slides: {len(prs.slides)}")
    print("All elements are fully editable in PowerPoint!")
    print("Ready for KLA interview!")

if __name__ == "__main__":
    main() 
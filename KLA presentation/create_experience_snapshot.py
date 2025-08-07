#!/usr/bin/env python3
"""
Create Experience Snapshot diagram/chart aligning resume experience with KLA SensArray role
Professional visual showing relevant experience categories
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_experience_snapshot_visual():
    """Create comprehensive experience snapshot visualization"""
    
    fig, ax = plt.subplots(figsize=(16, 12))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 12)
    ax.axis('off')
    
    # Set clean background
    fig.patch.set_facecolor('white')
    
    # Title
    ax.text(8, 11.5, 'EXPERIENCE SNAPSHOT', ha='center', va='center',
           fontsize=24, fontweight='bold', color='#1f4e79')
    ax.text(8, 11, 'Aligned with KLA SensArray Role Requirements', ha='center', va='center',
           fontsize=16, style='italic', color='#666666')
    
    # Experience categories with detailed content
    experiences = [
        {
            'title': 'Thin-Film Deposition\n(PVD, CVD, ALD)',
            'items': [
                '• Synthesized nano-thin films using spin coating',
                '  and sol-gel at Rayn Innovation',
                '• Conducted DOE with JMP to optimize deposition',
                '  rate and uniformity',
                '• Supported vacuum chamber setup and tool',
                '  troubleshooting at TSMC fabs'
            ],
            'position': (2, 8.5),
            'size': (6, 2.5),
            'color': '#4472c4',
            'icon': '🔬'
        },
        {
            'title': 'Vacuum System Design\n& Integration',
            'items': [
                '• Contributed to N+1 vacuum redundancy plan',
                '  at TSMC, reducing downtime by 25%',
                '• Routed vacuum utilities and coordinated',
                '  Fab21 tool hookups and layout',
                '• Experience with chamber sealing, leak testing,',
                '  and pressure regulation'
            ],
            'position': (10, 8.5),
            'size': (6, 2.5),
            'color': '#70ad47',
            'icon': '⚙️'
        },
        {
            'title': 'Root Cause Analysis\n8D / DOE',
            'items': [
                '• Applied JMP to DOE for thermal and film',
                '  thickness tuning at Rayn',
                '• Led FMEA and 8D to resolve vacuum chamber',
                '  leak at TSMC, saving $120K annually',
                '• Hands-on use of failure data to refine',
                '  tool calibration sequences'
            ],
            'position': (2, 5),
            'size': (6, 2.5),
            'color': '#ffc000',
            'icon': '🔍'
        },
        {
            'title': 'System Design, CAD\n& Thermal Analysis',
            'items': [
                '• Modeled system testbeds in SolidWorks',
                '  and Fusion360',
                '• Performed CFD for thermal cooling zones',
                '  in data center and pod design',
                '• Created engineering drawings for vacuum',
                '  piping and tool layout (TSMC/Chemtech)'
            ],
            'position': (10, 5),
            'size': (6, 2.5),
            'color': '#c5504b',
            'icon': '📐'
        }
    ]
    
    # Draw experience boxes
    for exp in experiences:
        x, y = exp['position']
        w, h = exp['size']
        
        # Main background box
        main_box = FancyBboxPatch((x-w/2, y-h/2), w, h, 
                                 boxstyle="round,pad=0.15", 
                                 facecolor=exp['color'], alpha=0.1,
                                 edgecolor=exp['color'], linewidth=3)
        ax.add_patch(main_box)
        
        # Header box
        header_box = FancyBboxPatch((x-w/2+0.1, y+h/2-0.6), w-0.2, 0.5,
                                   boxstyle="round,pad=0.1",
                                   facecolor=exp['color'], alpha=0.9)
        ax.add_patch(header_box)
        
        # Icon and title
        ax.text(x-w/2+0.5, y+h/2-0.35, exp['icon'], ha='center', va='center',
               fontsize=16, color='white')
        ax.text(x, y+h/2-0.35, exp['title'], ha='center', va='center',
               fontsize=12, fontweight='bold', color='white')
        
        # Experience items
        item_y = y + 0.7
        for item in exp['items']:
            ax.text(x-w/2+0.3, item_y, item, ha='left', va='top',
                   fontsize=10, color=exp['color'], fontweight='semibold')
            item_y -= 0.25
    
    # Add connecting lines to show integration
    # Connect related experiences
    connections = [
        ((5, 7.25), (11, 7.25)),  # Deposition to Vacuum
        ((5, 6.25), (11, 6.25)),  # RCA to CAD
        ((8, 7.25), (8, 6.25))    # Vertical connection
    ]
    
    for start, end in connections:
        ax.plot([start[0], end[0]], [start[1], end[1]], 
               color='#cccccc', linewidth=2, alpha=0.7, linestyle='--')
    
    # Add KLA relevance section at bottom
    kla_box = FancyBboxPatch((1, 0.5), 14, 1.5, 
                            boxstyle="round,pad=0.2",
                            facecolor='#1f4e79', alpha=0.1,
                            edgecolor='#1f4e79', linewidth=2)
    ax.add_patch(kla_box)
    
    ax.text(8, 1.7, 'KLA SENSARRAY ALIGNMENT', ha='center', va='center',
           fontsize=14, fontweight='bold', color='#1f4e79')
    
    alignment_points = [
        '✓ In-Situ Process Monitoring & Control',
        '✓ Semiconductor Deposition Tool Experience', 
        '✓ Vacuum System Integration & Design',
        '✓ Statistical Process Control & DOE',
        '✓ Root Cause Analysis & Problem Solving',
        '✓ CAD Design & Thermal Analysis'
    ]
    
    # Two columns of alignment points
    left_points = alignment_points[:3]
    right_points = alignment_points[3:]
    
    point_y = 1.4
    for point in left_points:
        ax.text(2.5, point_y, point, ha='left', va='center',
               fontsize=11, color='#1f4e79', fontweight='semibold')
        point_y -= 0.25
    
    point_y = 1.4
    for point in right_points:
        ax.text(8.5, point_y, point, ha='left', va='center',
               fontsize=11, color='#1f4e79', fontweight='semibold')
        point_y -= 0.25
    
    plt.tight_layout()
    plt.savefig('experience_snapshot.png', dpi=300, bbox_inches='tight',
               facecolor='white', edgecolor='none', pad_inches=0.3)
    plt.close()
    
    print("✅ Created experience snapshot: experience_snapshot.png")
    return 'experience_snapshot.png'

def create_experience_timeline():
    """Create timeline view of relevant experience"""
    
    fig, ax = plt.subplots(figsize=(14, 8))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    ax.axis('off')
    
    # Title
    ax.text(7, 7.5, 'RELEVANT EXPERIENCE TIMELINE', ha='center', va='center',
           fontsize=18, fontweight='bold', color='#1f4e79')
    
    # Timeline data
    timeline_items = [
        {
            'period': 'Jul 2024 – May 2025',
            'company': 'TSMC',
            'role': 'Mechanical Engineer',
            'highlights': ['Vacuum system design', 'Tool troubleshooting', '$120K cost savings'],
            'x_pos': 2.5,
            'color': '#4472c4'
        },
        {
            'period': 'Jan 2024 – May 2024',
            'company': 'Rayn Innovation',
            'role': 'Senior Opto-Mechanical Engineer',
            'highlights': ['Thin-film deposition', 'DOE optimization', 'JMP analysis'],
            'x_pos': 7,
            'color': '#70ad47'
        },
        {
            'period': 'Jun 2019 – May 2022',
            'company': 'Chemtech Systems',
            'role': 'Senior Mechanical Engineer',
            'highlights': ['FMEA methodology', 'CAD design', 'System optimization'],
            'x_pos': 11.5,
            'color': '#c5504b'
        }
    ]
    
    # Draw timeline line
    ax.plot([1, 13], [4, 4], color='#333333', linewidth=3)
    
    # Add timeline items
    for item in timeline_items:
        x = item['x_pos']
        
        # Timeline marker
        marker = Circle((x, 4), 0.2, facecolor=item['color'], 
                       edgecolor='white', linewidth=2, zorder=3)
        ax.add_patch(marker)
        
        # Company box
        company_box = FancyBboxPatch((x-1.2, 5), 2.4, 0.8,
                                    boxstyle="round,pad=0.1",
                                    facecolor=item['color'], alpha=0.9)
        ax.add_patch(company_box)
        
        # Company and role
        ax.text(x, 5.55, item['company'], ha='center', va='center',
               fontsize=12, fontweight='bold', color='white')
        ax.text(x, 5.25, item['role'], ha='center', va='center',
               fontsize=10, color='white')
        
        # Period
        ax.text(x, 3.5, item['period'], ha='center', va='center',
               fontsize=10, fontweight='bold', color='#333333')
        
        # Highlights
        highlight_y = 2.8
        for highlight in item['highlights']:
            ax.text(x, highlight_y, f'• {highlight}', ha='center', va='center',
                   fontsize=9, color=item['color'], fontweight='semibold')
            highlight_y -= 0.3
        
        # Connection line
        ax.plot([x, x], [4.2, 4.9], color=item['color'], linewidth=2)
    
    # Add legend
    legend_y = 0.8
    ax.text(7, legend_y, 'Key: Direct semiconductor manufacturing experience with deposition, vacuum systems, and process optimization',
           ha='center', va='center', fontsize=11, style='italic', color='#666666')
    
    plt.tight_layout()
    plt.savefig('experience_timeline.png', dpi=300, bbox_inches='tight',
               facecolor='white', edgecolor='none')
    plt.close()
    
    print("✅ Created experience timeline: experience_timeline.png")
    return 'experience_timeline.png'

def create_experience_powerpoint():
    """Create PowerPoint with experience visuals"""
    
    # Generate visuals
    snapshot_file = create_experience_snapshot_visual()
    timeline_file = create_experience_timeline()
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Experience Snapshot
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box1 = slide1.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(0.6))
    title_frame1 = title_box1.text_frame
    title_frame1.text = "Experience Snapshot - KLA SensArray Alignment"
    title_frame1.paragraphs[0].font.size = Pt(24)
    title_frame1.paragraphs[0].font.bold = True
    title_frame1.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    title_frame1.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add snapshot image
    slide1.shapes.add_picture(snapshot_file, Inches(0.5), Inches(0.9), 
                             width=Inches(9), height=Inches(6.8))
    
    # Slide 2: Experience Timeline
    slide2 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(0.6))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Relevant Experience Timeline"
    title_frame2.paragraphs[0].font.size = Pt(24)
    title_frame2.paragraphs[0].font.bold = True
    title_frame2.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    title_frame2.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add timeline image
    slide2.shapes.add_picture(timeline_file, Inches(0.5), Inches(0.9), 
                             width=Inches(9), height=Inches(6.8))
    
    # Save presentation
    prs.save('Experience_Snapshot_KLA.pptx')
    
    print("✅ Created experience presentation: Experience_Snapshot_KLA.pptx")
    
    return snapshot_file, timeline_file

def create_experience_guide():
    """Create presentation guide for experience snapshot"""
    
    guide_content = """EXPERIENCE SNAPSHOT - PRESENTATION GUIDE
======================================

PURPOSE:
Demonstrate direct alignment between your background and KLA SensArray role requirements

WHEN TO USE:
• As backup slide during Q&A about your experience
• If interviewer asks about specific technical background
• To reinforce your qualifications during role-fit discussion
• As supporting material for "Why you're the right candidate"

VISUAL ELEMENTS:

📊 FOUR EXPERIENCE CATEGORIES:

1. THIN-FILM DEPOSITION (Blue - Top Left)
   - Rayn Innovation: Nano-thin films, spin coating, sol-gel
   - DOE optimization with JMP
   - TSMC vacuum chamber support

2. VACUUM SYSTEM DESIGN (Green - Top Right)
   - TSMC N+1 redundancy plan (25% downtime reduction)
   - Fab21 tool hookups and layout
   - Chamber sealing and leak testing

3. ROOT CAUSE ANALYSIS (Yellow - Bottom Left)
   - JMP-based DOE at Rayn
   - 8D methodology at TSMC ($120K savings)
   - Failure data analysis for calibration

4. SYSTEM DESIGN & CAD (Red - Bottom Right)
   - SolidWorks and Fusion360 modeling
   - CFD thermal analysis
   - Engineering drawings for vacuum systems

🗣️ HOW TO PRESENT (2-3 minutes):

OPENING (30 seconds):
"Let me show you how my experience directly aligns with what KLA SensArray needs. I've organized my background into four key areas."

CATEGORY WALKTHROUGH (90 seconds):
"Starting with thin-film deposition - I have hands-on experience with the exact processes KLA's customers use. At Rayn Innovation, I synthesized nano-films and optimized deposition rates using DOE.

For vacuum systems, at TSMC I designed redundancy plans and coordinated tool installations - critical for the fab environment where SensArray operates.

My root cause analysis experience includes both the statistical rigor (JMP, DOE) and the practical problem-solving (8D, FMEA) that KLA values.

Finally, my CAD and thermal analysis skills directly support sensor wafer design and integration."

CLOSING (30 seconds):
"This isn't just related experience - this is direct, hands-on work in the exact technical domains where KLA SensArray operates."

KEY MESSAGES:
✅ Direct semiconductor manufacturing experience
✅ Hands-on deposition and vacuum system work
✅ Statistical process optimization expertise
✅ Proven cost savings and performance improvements
✅ Cross-functional technical skills

TECHNICAL DEPTH READY:
- Specific JMP techniques used
- Vacuum system design principles
- 8D methodology application
- CFD analysis approaches
- CAD modeling best practices

STRATEGIC VALUE:
• Proves you can contribute immediately
• Shows understanding of KLA's technical challenges
• Demonstrates both breadth and depth of experience
• Connects your background to business impact
"""
    
    with open('Experience_Snapshot_Guide.txt', 'w') as f:
        f.write(guide_content)
    
    print("✅ Created experience guide: Experience_Snapshot_Guide.txt")

if __name__ == "__main__":
    print("🎯 Creating Experience Snapshot visualization...")
    create_experience_powerpoint()
    create_experience_guide()
    print("\n✅ Experience Snapshot complete! Perfect backup material for your interview.")
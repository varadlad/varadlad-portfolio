#!/usr/bin/env python3
"""
KLA Manufacturing Design Engineer Presentation Generator
Creates a comprehensive PowerPoint presentation for Varad Lad's interview
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Set style for charts
plt.style.use('seaborn-v0_8')
sns.set_palette("husl")

def create_title_slide(prs):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Advanced Thermal & Vacuum Systems Engineering"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Varad Lad\nManufacturing Design Engineer Candidate\nKLA Corporation - SensArray Division"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[1].font.size = Pt(18)
    subtitle.text_frame.paragraphs[2].font.size = Pt(18)
    
    return slide

def create_intro_slide(prs):
    """Create introduction slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Engineering Excellence Across Semiconductor & Data Center Domains"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content = slide.placeholders[1]
    content_text = """• 5+ Years Experience: Mechanical Design, CAD Modeling, FEA/CFD Analysis
• Semiconductor Expertise: TSMC 4nm/2nm Fab Design, Thin-Film Deposition (CVD/PVD/ALD)
• Thermal Systems: Data Center Cooling Optimization, Vacuum Systems, HVAC Design
• Process Optimization: DOE, SPC, FMEA, Yield Improvement, Cost Reduction
• Technical Leadership: Cross-functional Teams, Stakeholder Management, Training
• Education: MS Mechanical Engineering, ASU Sun Award of Excellence"""
    
    content.text = content_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_main_project_slide(prs):
    """Create main project slide - Data Center Thermal Optimization"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Data Center Advanced Thermal Optimization"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content = slide.placeholders[1]
    content_text = """CHALLENGE: Data center cooling consumes 50% of total energy with thermal stratification issues

SOLUTION: Comprehensive CFD analysis comparing air vs. liquid cooling strategies
• 3D Modeling: Fusion 360 rack design with modular cooling features
• CFD Simulation: SimScale analysis under realistic 10kW/rack loads
• Hot Aisle Containment: Optimized airflow management
• Adaptive Control: Smart thermal management strategies

RESULTS:
• 27% improvement in cooling performance
• 32% reduction in hotspots
• 15°C cooler inlet temperatures with liquid cooling
• PUE improvement from 1.5 to 1.1 (liquid-cooled hybrid)
• Scalable, cost-effective thermal management approach"""
    
    content.text = content_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_thermal_analysis_slide(prs):
    """Create thermal analysis results slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Thermal Performance Analysis: Air vs. Liquid Cooling"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create performance comparison chart
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 6))
    
    # Temperature comparison
    categories = ['Air Cooling', 'Liquid Cooling']
    max_temp = [45, 30]
    avg_temp = [38, 25]
    
    x = np.arange(len(categories))
    width = 0.35
    
    ax1.bar(x - width/2, max_temp, width, label='Max Temperature (°C)', color='#ff6b6b', alpha=0.8)
    ax1.bar(x + width/2, avg_temp, width, label='Avg Temperature (°C)', color='#4ecdc4', alpha=0.8)
    
    ax1.set_ylabel('Temperature (°C)')
    ax1.set_title('Temperature Comparison')
    ax1.set_xticks(x)
    ax1.set_xticklabels(categories)
    ax1.legend()
    ax1.grid(True, alpha=0.3)
    
    # Performance metrics
    metrics = ['Cooling Performance', 'Hotspot Reduction', 'Energy Efficiency', 'Cost Effectiveness']
    air_scores = [73, 68, 67, 75]
    liquid_scores = [100, 100, 89, 85]
    
    x = np.arange(len(metrics))
    width = 0.35
    
    ax2.bar(x - width/2, air_scores, width, label='Air Cooling', color='#ff6b6b', alpha=0.8)
    ax2.bar(x + width/2, liquid_scores, width, label='Liquid Cooling', color='#4ecdc4', alpha=0.8)
    
    ax2.set_ylabel('Performance Score (%)')
    ax2.set_title('Performance Metrics Comparison')
    ax2.set_xticks(x)
    ax2.set_xticklabels(metrics, rotation=45, ha='right')
    ax2.legend()
    ax2.grid(True, alpha=0.3)
    
    plt.tight_layout()
    plt.savefig('thermal_analysis.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add image to slide
    slide.shapes.add_picture('thermal_analysis.png', Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    
    return slide

def create_thin_film_slide(prs):
    """Create thin film deposition slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Thin-Film Deposition Process Optimization"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content = slide.placeholders[1]
    content_text = """TECHNICAL EXPERTISE: CVD, PVD, ALD Deposition Techniques

PROCESS OPTIMIZATION:
• DOE & SPC Analysis: JMP-based parameter optimization
• Metrology: SEM, TEM, XRD, UV-Vis characterization
• Defect Detection: 90% surface defect identification
• Film Uniformity: 12% improvement through parameter tuning
• Material Waste: 40% reduction via process optimization

YIELD IMPROVEMENT:
• 28% yield performance gain through thermal management
• 15% synthesis-stage purity improvement
• Big-data analysis of deposition variables and defect trends
• Predictive modeling for process optimization

RELEVANCE TO KLA SENSARRAY:
• In-situ thermal process control expertise
• Vacuum system optimization experience
• Process monitoring and control systems
• Yield improvement methodologies"""
    
    content.text = content_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_vacuum_systems_slide(prs):
    """Create vacuum systems slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "N+1 Redundancy in Vacuum Systems (TSMC Experience)"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content = slide.placeholders[1]
    content_text = """CRITICAL FACILITY SYSTEM DESIGN:
• 4nm, 2nm Fab-Utility and Data Center Design
• N+1 Redundancy Planning with Operations Reliability Team
• $120,000 Future Maintenance Cost Savings
• 25% Equipment Downtime Reduction

VACUUM SYSTEM EXPERTISE:
• SCADA-based Diagnostics for Critical HVAC Systems
• Abatement System Failure Resolution
• Cleanroom Environment System Troubleshooting
• Cooling Towers and Chillers Optimization

QUALITY ASSURANCE:
• 0% Construction Rework Achievement
• P&ID and PFD Validation
• ICC, IBC, IFC, NFPA Code Compliance
• Power BI and Tableau Dashboard Development

RELEVANCE TO KLA:
• High Vacuum Systems Experience
• Process Control and Monitoring
• Equipment Reliability and Maintenance
• Cross-functional Team Leadership"""
    
    content.text = content_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_kla_alignment_slide(prs):
    """Create KLA SensArray alignment slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Alignment with KLA SensArray: In-Situ Thermal Process Control"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create alignment diagram
    fig, ax = plt.subplots(figsize=(12, 8))
    
    # Define categories and skills
    categories = ['Thermal Process\nControl', 'Vacuum Systems', 'Process\nOptimization', 'Equipment\nReliability', 'Yield\nImprovement']
    kla_needs = [95, 90, 85, 88, 92]
    varad_experience = [95, 88, 90, 85, 90]
    
    x = np.arange(len(categories))
    width = 0.35
    
    bars1 = ax.bar(x - width/2, kla_needs, width, label='KLA SensArray Requirements', 
                   color='#2E86AB', alpha=0.8)
    bars2 = ax.bar(x + width/2, varad_experience, width, label='Varad\'s Experience', 
                   color='#A23B72', alpha=0.8)
    
    ax.set_ylabel('Expertise Level (%)')
    ax.set_title('Skills Alignment Matrix')
    ax.set_xticks(x)
    ax.set_xticklabels(categories)
    ax.legend()
    ax.grid(True, alpha=0.3)
    ax.set_ylim(0, 100)
    
    # Add value labels on bars
    for bar in bars1:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height + 1,
                f'{int(height)}%', ha='center', va='bottom', fontweight='bold')
    
    for bar in bars2:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height + 1,
                f'{int(height)}%', ha='center', va='bottom', fontweight='bold')
    
    plt.tight_layout()
    plt.savefig('kla_alignment.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add image to slide
    slide.shapes.add_picture('kla_alignment.png', Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    
    return slide

def create_technical_skills_slide(prs):
    """Create technical skills slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Technical Skills & Tools"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content = slide.placeholders[1]
    content_text = """CAD DESIGN & SIMULATION:
• Revit, Siemens NX, ANSYS CFD/FEA (Icepak/Fluent)
• AutoCAD, SolidWorks, CATIA, Creo, Fusion 360, SimScale

DATA ANALYSIS & CODING:
• Python, MATLAB, JMP, SQL, DoE, SPC
• Power BI, Tableau, RCA, FMEA, SAP HANA
• PLC, HMI, G-code

THERMAL & VACUUM SYSTEMS:
• HVAC Design and Optimization
• Vacuum Deposition Equipment (CVD, PVD, ALD)
• Thermal Management and Heat Transfer
• Process Control and Automation

LEADERSHIP & COMMUNICATION:
• President of Employee Collaboration Club at TSMC
• Technical Public Speaker and Mentor
• Cross-functional Team Leadership
• Stakeholder Management and Training"""
    
    content.text = content_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_achievements_slide(prs):
    """Create achievements and metrics slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Quantified Impact & Achievements"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Create achievements visualization
    fig, ax = plt.subplots(figsize=(12, 8))
    
    # Define achievement categories and values
    categories = ['Cost Savings\n($120K)', 'Downtime\nReduction\n(25%)', 'Yield\nImprovement\n(28%)', 'Film\nUniformity\n(12%)', 'Material\nWaste\nReduction\n(40%)']
    values = [120, 25, 28, 12, 40]
    colors = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6B5B95']
    
    # Create horizontal bar chart
    bars = ax.barh(categories, values, color=colors, alpha=0.8)
    
    # Add value labels
    for i, (bar, value) in enumerate(zip(bars, values)):
        ax.text(bar.get_width() + 2, bar.get_y() + bar.get_height()/2,
                f'{value}%' if i > 0 else f'${value}K', 
                va='center', fontweight='bold', fontsize=12)
    
    ax.set_xlabel('Improvement Percentage / Cost Savings')
    ax.set_title('Key Performance Metrics')
    ax.grid(True, alpha=0.3)
    ax.set_xlim(0, max(values) * 1.2)
    
    plt.tight_layout()
    plt.savefig('achievements.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add image to slide
    slide.shapes.add_picture('achievements.png', Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    
    return slide

def create_qa_slide(prs):
    """Create Q&A slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Thank You & Questions"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content = slide.placeholders[1]
    content_text = """KEY VALUE PROPOSITION:
• Proven track record in thermal process control and vacuum systems
• Experience with semiconductor manufacturing and yield improvement
• Strong background in equipment reliability and process optimization
• Demonstrated leadership in cross-functional engineering teams

CONTACT INFORMATION:
• Email: vlad3@asu.edu
• Phone: 602-388-6861
• LinkedIn: linkedin.com/in/varadlad
• Portfolio: varadlad.github.io/varadlad-portfolio

AVAILABILITY:
• Available immediately
• Open to relocation
• Ready to contribute to KLA's innovative mission

Thank you for your time and consideration!"""
    
    content.text = content_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def main():
    """Main function to create the presentation"""
    print("Creating KLA Manufacturing Design Engineer Presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create slides
    create_title_slide(prs)
    create_intro_slide(prs)
    create_main_project_slide(prs)
    create_thermal_analysis_slide(prs)
    create_thin_film_slide(prs)
    create_vacuum_systems_slide(prs)
    create_kla_alignment_slide(prs)
    create_technical_skills_slide(prs)
    create_achievements_slide(prs)
    create_qa_slide(prs)
    
    # Save presentation
    output_file = "Varad_Lad_KLA_Presentation.pptx"
    prs.save(output_file)
    
    # Clean up temporary files
    for temp_file in ['thermal_analysis.png', 'kla_alignment.png', 'achievements.png']:
        if os.path.exists(temp_file):
            os.remove(temp_file)
    
    print(f"Presentation created successfully: {output_file}")
    print("Total slides: 10")
    print("Ready for KLA interview!")

if __name__ == "__main__":
    main() 
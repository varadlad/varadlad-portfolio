#!/usr/bin/env python3
"""
Create clean radial cluster diagram for Slide 4: User Requirements
Single diagram with no overlapping text and professional layout
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Circle, Wedge, Arrow
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_clean_radial_diagram():
    """Create single, clean radial cluster diagram with no overlapping text"""
    
    # Create figure with larger size for better spacing
    fig, ax = plt.subplots(figsize=(14, 14))
    ax.set_xlim(-12, 12)
    ax.set_ylim(-12, 12)
    ax.set_aspect('equal')
    ax.axis('off')
    
    # Set clean white background
    fig.patch.set_facecolor('white')
    
    # Center circle - Main objective (larger and more prominent)
    center_circle = Circle((0, 0), 3, facecolor='#1f4e79', edgecolor='white', linewidth=4, zorder=5)
    ax.add_patch(center_circle)
    
    # Center text with better spacing
    ax.text(0, 0.8, 'DEPOSITION', ha='center', va='center', fontsize=18, fontweight='bold', color='white', zorder=6)
    ax.text(0, 0, 'RATE', ha='center', va='center', fontsize=18, fontweight='bold', color='white', zorder=6)
    ax.text(0, -0.8, 'OPTIMIZATION', ha='center', va='center', fontsize=18, fontweight='bold', color='white', zorder=6)
    
    # Requirements with carefully calculated positions to avoid overlap
    requirements = [
        {
            'title': 'Higher\nThroughput',
            'metric': '+40%',
            'detail': 'Boost manufacturing\nproductivity',
            'angle': 90,  # Top
            'color': '#4472c4',
            'text_color': 'white'
        },
        {
            'title': 'Quality\nPreservation', 
            'metric': 'Maintain\nSpecs',
            'detail': 'Film thickness &\nperformance standards',
            'angle': 18,  # Top-right
            'color': '#70ad47',
            'text_color': 'white'
        },
        {
            'title': 'Process\nConsistency',
            'metric': '-25%\nVariability',
            'detail': 'Reduce defects &\nprocess variation',
            'angle': 306,  # Bottom-right
            'color': '#ffc000',
            'text_color': 'black'
        },
        {
            'title': 'Cost\nEfficiency',
            'metric': '-15%\nPer Wafer',
            'detail': 'Lower operating\ncosts & waste',
            'angle': 234,  # Bottom-left
            'color': '#c5504b',
            'text_color': 'white'
        },
        {
            'title': 'Operational\nConstraints',
            'metric': 'Within\nLimits',
            'detail': 'Equipment safety &\nfab compatibility',
            'angle': 162,  # Top-left
            'color': '#7030a0',
            'text_color': 'white'
        }
    ]
    
    # Draw requirements with better spacing
    for req in requirements:
        angle_rad = np.radians(req['angle'])
        
        # Calculate positions with more spacing
        spoke_start_x = 3.2 * np.cos(angle_rad)
        spoke_start_y = 3.2 * np.sin(angle_rad)
        spoke_end_x = 6.8 * np.cos(angle_rad)
        spoke_end_y = 6.8 * np.sin(angle_rad)
        circle_x = 7.5 * np.cos(angle_rad)
        circle_y = 7.5 * np.sin(angle_rad)
        
        # Draw spoke line with gradient effect
        ax.plot([spoke_start_x, spoke_end_x], 
               [spoke_start_y, spoke_end_y], 
               color=req['color'], linewidth=5, alpha=0.7, zorder=2)
        
        # Draw requirement circle (larger for better text fit)
        req_circle = Circle((circle_x, circle_y), 2.2, 
                           facecolor=req['color'], edgecolor='white', 
                           linewidth=3, alpha=0.95, zorder=3)
        ax.add_patch(req_circle)
        
        # Add title text (better positioning)
        ax.text(circle_x, circle_y + 0.7, req['title'], ha='center', va='center',
               fontsize=12, fontweight='bold', color=req['text_color'], zorder=4)
        
        # Add metric text
        ax.text(circle_x, circle_y - 0.2, req['metric'], ha='center', va='center',
               fontsize=11, fontweight='bold', color=req['text_color'], zorder=4)
        
        # Position detail text boxes outside circles with smart positioning
        detail_distance = 10.5
        detail_x = detail_distance * np.cos(angle_rad)
        detail_y = detail_distance * np.sin(angle_rad)
        
        # Adjust text box positioning to avoid overlaps
        box_width = 2.8
        box_height = 1.2
        
        # Smart positioning based on angle
        if req['angle'] < 90 or req['angle'] > 270:  # Right side
            text_align = 'left'
            box_x = detail_x - 0.2
        else:  # Left side
            text_align = 'right'
            box_x = detail_x - box_width + 0.2
            
        if 45 < req['angle'] < 135:  # Top
            box_y = detail_y - 0.2
        elif 225 < req['angle'] < 315:  # Bottom
            box_y = detail_y - box_height + 0.2
        else:  # Sides
            box_y = detail_y - box_height/2
        
        # Detail text background box
        detail_box = FancyBboxPatch((box_x, box_y), box_width, box_height,
                                   boxstyle="round,pad=0.15", 
                                   facecolor='white', edgecolor=req['color'],
                                   linewidth=2, alpha=0.95, zorder=1)
        ax.add_patch(detail_box)
        
        # Detail text
        ax.text(box_x + box_width/2, box_y + box_height/2, req['detail'], 
               ha='center', va='center', fontsize=10, color=req['color'], 
               fontweight='semibold', zorder=2)
        
        # Connection line from circle to detail box
        connection_start_x = circle_x + 2.2 * np.cos(angle_rad)
        connection_start_y = circle_y + 2.2 * np.sin(angle_rad)
        connection_end_x = box_x + box_width/2
        connection_end_y = box_y + box_height/2
        
        ax.plot([connection_start_x, connection_end_x], 
               [connection_start_y, connection_end_y], 
               color=req['color'], linewidth=1.5, alpha=0.6, 
               linestyle='--', zorder=1)
    
    # Add title with better positioning
    ax.text(0, 11, 'USER REQUIREMENTS', ha='center', va='center',
           fontsize=24, fontweight='bold', color='#1f4e79')
    ax.text(0, 10.2, 'Integrated Optimization Goals', ha='center', va='center',
           fontsize=16, style='italic', color='#666666')
    
    # Add subtle background grid for professional look
    circle_bg = Circle((0, 0), 9, facecolor='none', edgecolor='#f0f0f0', 
                      linewidth=1, alpha=0.5, zorder=0)
    ax.add_patch(circle_bg)
    
    # Save with high quality
    plt.tight_layout()
    plt.savefig('slide4_clean_radial_diagram.png', dpi=300, bbox_inches='tight',
               facecolor='white', edgecolor='none', pad_inches=0.2)
    plt.close()
    
    print("✅ Created clean radial diagram: slide4_clean_radial_diagram.png")
    return 'slide4_clean_radial_diagram.png'

def create_slide4_powerpoint_clean():
    """Create clean PowerPoint slide with the radial diagram"""
    
    # Generate the clean diagram
    diagram_file = create_clean_radial_diagram()
    
    # Create new PowerPoint presentation
    prs = Presentation()
    
    # Create slide with blank layout
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "User Requirements"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)  # #1f4e79
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add the clean radial diagram (centered and properly sized)
    slide.shapes.add_picture(diagram_file, Inches(0.5), Inches(1), 
                           width=Inches(9), height=Inches(7))
    
    # Save presentation
    prs.save('Slide4_Clean_Requirements.pptx')
    
    print("✅ Created clean PowerPoint: Slide4_Clean_Requirements.pptx")
    
    return diagram_file

def cleanup_old_files():
    """Remove old diagram files"""
    import os
    
    old_files = [
        'user_requirements_radial_diagram.png',
        'user_requirements_cycle_diagram.png',
        'Slide4_Visual_Requirements.pptx'
    ]
    
    for file in old_files:
        try:
            if os.path.exists(file):
                os.remove(file)
                print(f"🗑️ Removed old file: {file}")
        except Exception as e:
            print(f"⚠️ Could not remove {file}: {e}")

def create_presentation_guide():
    """Create guide for the clean radial diagram"""
    
    guide_content = """SLIDE 4 - CLEAN RADIAL REQUIREMENTS DIAGRAM
==========================================

DIAGRAM FEATURES:
✅ Single, professional radial cluster design
✅ No overlapping text - all elements clearly spaced
✅ Central hub with main objective
✅ 5 color-coded requirement clusters
✅ Quantified metrics prominently displayed
✅ Detailed descriptions in separate boxes
✅ Clean white background for projection clarity

VISUAL ELEMENTS:
🎯 CENTER: "Deposition Rate Optimization" - main objective
📊 SPOKES: 5 requirements radiating outward
🎨 COLORS: Professional color scheme (blue, green, yellow, red, purple)
📋 METRICS: Specific targets clearly visible
📝 DETAILS: Supporting information in connected boxes

HOW TO PRESENT:

🗣️ OPENING (30 seconds):
"This radial diagram shows our five integrated optimization requirements. Rather than treating these as trade-offs, our approach aims to achieve all goals simultaneously."

🗣️ CLOCKWISE PRESENTATION (2-3 minutes):
1. TOP (Blue): "Higher Throughput - targeting 40% improvement in deposition rate"
2. TOP-RIGHT (Green): "Quality Preservation - maintaining all film specifications"
3. BOTTOM-RIGHT (Yellow): "Process Consistency - reducing variability by 25%"
4. BOTTOM-LEFT (Red): "Cost Efficiency - achieving 15% savings per wafer"
5. TOP-LEFT (Purple): "Operational Constraints - staying within equipment limits"

🗣️ CLOSING (30 seconds):
"The center represents our integrated solution - optimizing deposition rate while satisfying all requirements through systematic methodology."

KEY ADVANTAGES:
✅ Visual clarity - immediate comprehension
✅ Professional appearance - no text overload
✅ Balanced presentation - equal weight to all requirements
✅ Quantified targets - specific and measurable
✅ Technical credibility - systematic approach

TECHNICAL DEPTH READY:
- Why 40% improvement is achievable
- How quality is maintained during optimization
- Statistical basis for variability reduction targets
- Cost calculation methodology
- Equipment constraint definitions
"""
    
    with open('Slide4_Clean_Guide.txt', 'w') as f:
        f.write(guide_content)
    
    print("✅ Created presentation guide: Slide4_Clean_Guide.txt")

if __name__ == "__main__":
    print("🎯 Creating clean radial diagram for Slide 4...")
    cleanup_old_files()
    create_slide4_powerpoint_clean()
    create_presentation_guide()
    print("\n✅ Clean Slide 4 complete! Use 'Slide4_Clean_Requirements.pptx'")
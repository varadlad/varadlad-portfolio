#!/usr/bin/env python3
"""
Create visual diagram for Slide 4: User Requirements
Using radial/spoke design with minimal text and maximum visual impact
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Circle, Wedge
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_user_requirements_diagram():
    """Create radial spoke diagram for user requirements"""
    
    fig, ax = plt.subplots(figsize=(12, 12))
    ax.set_xlim(-10, 10)
    ax.set_ylim(-10, 10)
    ax.set_aspect('equal')
    ax.axis('off')
    
    # Center circle - Main objective
    center_circle = Circle((0, 0), 2.5, facecolor='#1f4e79', edgecolor='white', linewidth=3)
    ax.add_patch(center_circle)
    ax.text(0, 0.2, 'OPTIMIZE', ha='center', va='center', fontsize=16, fontweight='bold', color='white')
    ax.text(0, -0.2, 'DEPOSITION', ha='center', va='center', fontsize=16, fontweight='bold', color='white')
    ax.text(0, -0.8, 'PROCESS', ha='center', va='center', fontsize=16, fontweight='bold', color='white')
    
    # Requirements with their positions, colors, icons, and metrics
    requirements = [
        {
            'title': 'Higher\nThroughput',
            'metric': '30-40%\nIncrease',
            'angle': 0,
            'color': '#4472c4',
            'icon': '⚡',
            'description': 'Faster deposition\nrate for improved\nproductivity'
        },
        {
            'title': 'Quality\nPreservation', 
            'metric': 'Maintain\nSpecs',
            'angle': 72,
            'color': '#70ad47',
            'icon': '🎯',
            'description': 'Film thickness\nuniformity &\nperformance'
        },
        {
            'title': 'Process\nConsistency',
            'metric': '20-25%\nReduction',
            'angle': 144,
            'color': '#ffc000',
            'icon': '📊',
            'description': 'Lower variability\n& defect rates'
        },
        {
            'title': 'Cost\nEfficiency',
            'metric': '15%\nSavings',
            'angle': 216,
            'color': '#c5504b',
            'icon': '💰',
            'description': 'Lower operating\ncost per wafer'
        },
        {
            'title': 'Operational\nConstraints',
            'metric': 'Within\nLimits',
            'angle': 288,
            'color': '#7030a0',
            'icon': '⚙️',
            'description': 'Equipment safety\n& compatibility'
        }
    ]
    
    # Draw spokes and requirement circles
    for req in requirements:
        angle_rad = np.radians(req['angle'])
        
        # Calculate positions
        spoke_end_x = 6.5 * np.cos(angle_rad)
        spoke_end_y = 6.5 * np.sin(angle_rad)
        circle_x = 7.5 * np.cos(angle_rad)
        circle_y = 7.5 * np.sin(angle_rad)
        
        # Draw spoke line
        ax.plot([2.5 * np.cos(angle_rad), spoke_end_x], 
               [2.5 * np.sin(angle_rad), spoke_end_y], 
               color=req['color'], linewidth=4, alpha=0.8)
        
        # Draw requirement circle
        req_circle = Circle((circle_x, circle_y), 1.8, 
                           facecolor=req['color'], edgecolor='white', 
                           linewidth=2, alpha=0.9)
        ax.add_patch(req_circle)
        
        # Add icon (using text as emoji might not render well, use symbols)
        icon_symbols = {'⚡': '⚡', '🎯': '●', '📊': '▣', '💰': '$', '⚙️': '⚙'}
        symbol = icon_symbols.get(req['icon'], '●')
        ax.text(circle_x, circle_y + 0.7, symbol, ha='center', va='center', 
               fontsize=24, color='white', fontweight='bold')
        
        # Add title
        ax.text(circle_x, circle_y + 0.1, req['title'], ha='center', va='center',
               fontsize=11, fontweight='bold', color='white')
        
        # Add metric
        ax.text(circle_x, circle_y - 0.6, req['metric'], ha='center', va='center',
               fontsize=9, fontweight='bold', color='white')
        
        # Add description box outside
        desc_x = 9.2 * np.cos(angle_rad)
        desc_y = 9.2 * np.sin(angle_rad)
        
        # Description background
        desc_box = FancyBboxPatch((desc_x - 1, desc_y - 0.5), 2, 1,
                                 boxstyle="round,pad=0.1", 
                                 facecolor='white', edgecolor=req['color'],
                                 linewidth=1, alpha=0.9)
        ax.add_patch(desc_box)
        
        ax.text(desc_x, desc_y, req['description'], ha='center', va='center',
               fontsize=8, color=req['color'], fontweight='bold')
    
    # Add title
    ax.text(0, 9.5, 'USER REQUIREMENTS', ha='center', va='center',
           fontsize=20, fontweight='bold', color='#1f4e79')
    ax.text(0, 8.8, 'Deposition Rate Optimization Goals', ha='center', va='center',
           fontsize=14, style='italic', color='#666666')
    
    # Add connecting arrows between related requirements
    # Connect Quality to Consistency
    ax.annotate('', xy=(requirements[2]['angle'], 6), xytext=(requirements[1]['angle'], 6),
               arrowprops=dict(arrowstyle='->', lw=1.5, color='gray', alpha=0.5,
                             connectionstyle="arc3,rad=0.3"))
    
    plt.tight_layout()
    plt.savefig('user_requirements_radial_diagram.png', dpi=300, bbox_inches='tight',
               facecolor='white', edgecolor='none')
    plt.close()
    return 'user_requirements_radial_diagram.png'

def create_alternative_cycle_diagram():
    """Create alternative cycle diagram showing requirement relationships"""
    
    fig, ax = plt.subplots(figsize=(12, 10))
    ax.set_xlim(-8, 8)
    ax.set_ylim(-6, 6)
    ax.set_aspect('equal')
    ax.axis('off')
    
    # Define cycle positions (pentagon shape)
    requirements = [
        {'title': 'Higher\nThroughput', 'metric': '+40%', 'pos': (0, 4), 'color': '#4472c4'},
        {'title': 'Quality\nPreservation', 'metric': 'Maintain', 'pos': (4, 2), 'color': '#70ad47'},
        {'title': 'Process\nConsistency', 'metric': '-25%', 'pos': (3, -2), 'color': '#ffc000'},
        {'title': 'Cost\nEfficiency', 'metric': '-15%', 'pos': (-3, -2), 'color': '#c5504b'},
        {'title': 'Operational\nConstraints', 'metric': 'Safe', 'pos': (-4, 2), 'color': '#7030a0'}
    ]
    
    # Draw connecting lines between requirements (cycle)
    for i in range(len(requirements)):
        start_pos = requirements[i]['pos']
        end_pos = requirements[(i + 1) % len(requirements)]['pos']
        
        ax.plot([start_pos[0], end_pos[0]], [start_pos[1], end_pos[1]], 
               color='lightgray', linewidth=3, alpha=0.7, zorder=1)
    
    # Draw requirement nodes
    for req in requirements:
        x, y = req['pos']
        
        # Main circle
        circle = Circle((x, y), 1.2, facecolor=req['color'], 
                       edgecolor='white', linewidth=3, zorder=2)
        ax.add_patch(circle)
        
        # Title
        ax.text(x, y + 0.3, req['title'], ha='center', va='center',
               fontsize=10, fontweight='bold', color='white', zorder=3)
        
        # Metric
        ax.text(x, y - 0.3, req['metric'], ha='center', va='center',
               fontsize=12, fontweight='bold', color='white', zorder=3)
    
    # Center objective
    center_circle = Circle((0, 0), 1.5, facecolor='#1f4e79', 
                          edgecolor='white', linewidth=3, zorder=2)
    ax.add_patch(center_circle)
    ax.text(0, 0.2, 'OPTIMAL', ha='center', va='center',
           fontsize=12, fontweight='bold', color='white', zorder=3)
    ax.text(0, -0.2, 'DEPOSITION', ha='center', va='center',
           fontsize=12, fontweight='bold', color='white', zorder=3)
    
    # Title
    ax.text(0, 5.5, 'INTEGRATED REQUIREMENTS CYCLE', ha='center', va='center',
           fontsize=18, fontweight='bold', color='#1f4e79')
    
    plt.tight_layout()
    plt.savefig('user_requirements_cycle_diagram.png', dpi=300, bbox_inches='tight',
               facecolor='white', edgecolor='none')
    plt.close()
    return 'user_requirements_cycle_diagram.png'

def create_slide4_powerpoint():
    """Create PowerPoint slide with visual diagram"""
    
    # Generate the diagrams
    radial_file = create_user_requirements_diagram()
    cycle_file = create_alternative_cycle_diagram()
    
    # Create PowerPoint presentation
    prs = Presentation()
    
    # Slide 1: Radial Diagram
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "User Requirements"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add radial diagram
    slide.shapes.add_picture(radial_file, Inches(0.5), Inches(1.3), 
                           width=Inches(9), height=Inches(6.5))
    
    # Slide 2: Cycle Diagram Alternative
    slide2 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "User Requirements - Integrated Approach"
    title_frame2.paragraphs[0].font.size = Pt(28)
    title_frame2.paragraphs[0].font.bold = True
    title_frame2.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame2.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add cycle diagram
    slide2.shapes.add_picture(cycle_file, Inches(0.5), Inches(1.3), 
                            width=Inches(9), height=Inches(6.5))
    
    # Save presentation
    prs.save('Slide4_Visual_Requirements.pptx')
    
    print("✅ Created visual diagrams:")
    print(f"  📊 Radial diagram: {radial_file}")
    print(f"  🔄 Cycle diagram: {cycle_file}")
    print(f"  📑 PowerPoint: Slide4_Visual_Requirements.pptx")
    
    return radial_file, cycle_file

def create_requirements_guide():
    """Create guide for presenting the visual requirements"""
    
    guide_content = """
SLIDE 4 VISUAL REQUIREMENTS GUIDE
=================================

DIAGRAM EXPLANATION:

🎯 RADIAL SPOKE DIAGRAM (Primary Option):
- CENTER CIRCLE: "Optimize Deposition Process" - the main objective
- 5 SPOKES: Each represents a key requirement category
- COLOR CODING: Different colors for easy identification
- METRICS: Quantified targets for each requirement
- DESCRIPTIONS: Brief explanations outside each circle

📊 CYCLE DIAGRAM (Alternative Option):
- PENTAGON SHAPE: Shows interconnected requirements
- CENTER HUB: Optimal deposition process
- CONNECTING LINES: Shows how requirements relate to each other
- BALANCED APPROACH: All requirements equally important

HOW TO PRESENT:

🗣️ OPENING:
"Rather than list requirements as bullet points, this diagram shows our optimization goals visually. The center represents our main objective - optimizing the deposition process. Each spoke represents a critical requirement."

🗣️ WALK THROUGH EACH SPOKE:
1. "Higher Throughput (blue) - We're targeting 30-40% improvement in deposition rate"
2. "Quality Preservation (green) - Maintaining all film specifications while going faster"  
3. "Process Consistency (yellow) - Reducing variability and defects by 20-25%"
4. "Cost Efficiency (red) - Achieving 15% cost savings per wafer"
5. "Operational Constraints (purple) - Staying within equipment safety limits"

🗣️ KEY MESSAGE:
"Notice these aren't trade-offs - our goal is to achieve ALL these requirements simultaneously through intelligent optimization, not by sacrificing one for another."

VISUAL ADVANTAGES:
✅ Immediate comprehension - audience sees all requirements at once
✅ Balanced presentation - no hierarchy implied
✅ Quantified targets - specific metrics clearly displayed
✅ Professional appearance - reduces text-heavy slide syndrome
✅ Memorable format - easier to recall than bullet points

TECHNICAL DEPTH READY:
- Be prepared to explain WHY each target is realistic
- Know the relationships between requirements
- Understand which requirements might conflict and how optimization resolves this
- Connect each requirement to semiconductor manufacturing best practices
"""
    
    with open('Slide4_Visual_Guide.txt', 'w') as f:
        f.write(guide_content)
    
    print("✅ Created presentation guide: Slide4_Visual_Guide.txt")

if __name__ == "__main__":
    create_slide4_powerpoint()
    create_requirements_guide()
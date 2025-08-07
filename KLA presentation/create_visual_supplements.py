#!/usr/bin/env python3
"""
Visual Supplements Generator for Deposition Rate Optimization Presentation
Creates charts, diagrams, tables, and visual elements to complement the main presentation
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import pandas as pd
import seaborn as sns
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle, Arrow
from matplotlib.sankey import Sankey
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Set style for professional charts
plt.style.use('seaborn-v0_8-whitegrid')
sns.set_palette("husl")

def create_process_flow_diagram():
    """Create process flow diagram for methodology (Slide 5)"""
    fig, ax = plt.subplots(1, 1, figsize=(12, 8))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis('off')
    
    # Phase 1: DOE
    doe_box = FancyBboxPatch((0.5, 5.5), 3, 1.5, boxstyle="round,pad=0.1", 
                            facecolor='lightblue', edgecolor='navy', linewidth=2)
    ax.add_patch(doe_box)
    ax.text(2, 6.25, 'Phase 1:\nDesign of Experiments\n(DOE)', ha='center', va='center', 
            fontsize=12, fontweight='bold')
    
    # Phase 2: Bayesian Optimization
    bayes_box = FancyBboxPatch((6, 5.5), 3, 1.5, boxstyle="round,pad=0.1", 
                              facecolor='lightgreen', edgecolor='darkgreen', linewidth=2)
    ax.add_patch(bayes_box)
    ax.text(7.5, 6.25, 'Phase 2:\nBayesian\nOptimization', ha='center', va='center', 
            fontsize=12, fontweight='bold')
    
    # Data flow arrows
    ax.arrow(3.5, 6.25, 2, 0, head_width=0.2, head_length=0.3, fc='black', ec='black')
    ax.text(4.75, 6.8, 'Model Training Data', ha='center', fontsize=10)
    
    # Feedback loop
    ax.arrow(7.5, 5.3, 0, -2, head_width=0.2, head_length=0.2, fc='red', ec='red')
    ax.arrow(7.5, 3, -5, 0, head_width=0.2, head_length=0.2, fc='red', ec='red')
    ax.arrow(2, 3, 0, 2.3, head_width=0.2, head_length=0.2, fc='red', ec='red')
    ax.text(5, 2.5, 'Iterative Refinement Loop', ha='center', fontsize=10, color='red')
    
    # Sub-processes
    ax.text(2, 4.5, '• Parameter Selection\n• Factorial Design\n• Initial Experiments', 
            ha='center', va='top', fontsize=10, bbox=dict(boxstyle="round,pad=0.3", facecolor='white'))
    
    ax.text(7.5, 4.5, '• Gaussian Process\n• Acquisition Function\n• Experiment Selection', 
            ha='center', va='top', fontsize=10, bbox=dict(boxstyle="round,pad=0.3", facecolor='white'))
    
    plt.title('Two-Phase Optimization Methodology', fontsize=16, fontweight='bold', pad=20)
    plt.tight_layout()
    plt.savefig('process_flow_diagram.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'process_flow_diagram.png'

def create_parameter_ranges_table():
    """Create parameter ranges table (Slide 7)"""
    # Sample parameter data
    parameters_data = {
        'Parameter': ['Temperature (°C)', 'Pressure (Torr)', 'Plasma Power (W)', 
                     'Gas Flow (sccm)', 'Spin Speed (RPM)', 'Time (min)'],
        'Minimum': [200, 0.5, 50, 10, 1000, 2],
        'Maximum': [500, 5.0, 150, 100, 3000, 10],
        'Baseline': [300, 2.0, 100, 50, 2000, 5],
        'Optimized': [380, 1.8, 145, 75, 2500, 4.2],
        'Impact Level': ['High', 'Medium', 'High', 'Low', 'Medium', 'High']
    }
    
    df = pd.DataFrame(parameters_data)
    
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.axis('tight')
    ax.axis('off')
    
    table = ax.table(cellText=df.values, colLabels=df.columns, cellLoc='center', loc='center',
                    colWidths=[0.2, 0.15, 0.15, 0.15, 0.15, 0.2])
    
    # Format table
    table.auto_set_font_size(False)
    table.set_fontsize(11)
    table.scale(1.2, 1.5)
    
    # Color code impact levels
    for i in range(len(df)):
        impact = df.iloc[i]['Impact Level']
        if impact == 'High':
            table[(i+1, 5)].set_facecolor('#ffcccc')
        elif impact == 'Medium':
            table[(i+1, 5)].set_facecolor('#ffffcc')
        else:
            table[(i+1, 5)].set_facecolor('#ccffcc')
    
    # Header formatting
    for j in range(len(df.columns)):
        table[(0, j)].set_facecolor('#4472C4')
        table[(0, j)].set_text_props(weight='bold', color='white')
    
    plt.title('Process Parameter Ranges and Optimization Results', fontsize=14, fontweight='bold', pad=20)
    plt.savefig('parameter_ranges_table.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'parameter_ranges_table.png'

def create_bayesian_optimization_workflow():
    """Create Bayesian optimization workflow diagram (Slide 8)"""
    fig, ax = plt.subplots(1, 1, figsize=(12, 10))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis('off')
    
    # Workflow steps
    steps = [
        {'pos': (2, 8.5), 'text': '1. Initialize\nwith DOE data', 'color': 'lightblue'},
        {'pos': (5, 8.5), 'text': '2. Update GP\nModel', 'color': 'lightgreen'},
        {'pos': (8, 8.5), 'text': '3. Acquisition\nFunction', 'color': 'lightyellow'},
        {'pos': (8, 6), 'text': '4. Select Next\nExperiment', 'color': 'lightcoral'},
        {'pos': (5, 6), 'text': '5. Run\nExperiment', 'color': 'lightpink'},
        {'pos': (2, 6), 'text': '6. Measure\nResults', 'color': 'lightgray'},
        {'pos': (5, 3.5), 'text': 'Converged?', 'color': 'lightsalmon'},
        {'pos': (2, 1.5), 'text': 'Optimal\nParameters', 'color': 'lightsteelblue'}
    ]
    
    # Draw boxes and text
    for step in steps:
        if 'Converged' in step['text']:
            # Diamond shape for decision
            diamond = patches.RegularPolygon(step['pos'], 4, radius=0.8, 
                                           orientation=np.pi/4, facecolor=step['color'], 
                                           edgecolor='black', linewidth=2)
            ax.add_patch(diamond)
        else:
            box = FancyBboxPatch((step['pos'][0]-0.7, step['pos'][1]-0.5), 1.4, 1, 
                               boxstyle="round,pad=0.1", facecolor=step['color'], 
                               edgecolor='black', linewidth=2)
            ax.add_patch(box)
        
        ax.text(step['pos'][0], step['pos'][1], step['text'], ha='center', va='center', 
                fontsize=10, fontweight='bold')
    
    # Arrows
    arrows = [
        ((2.7, 8.5), (4.3, 8.5)),  # 1 to 2
        ((5.7, 8.5), (7.3, 8.5)),  # 2 to 3
        ((8, 7.8), (8, 6.7)),      # 3 to 4
        ((7.3, 6), (5.7, 6)),      # 4 to 5
        ((4.3, 6), (2.7, 6)),      # 5 to 6
        ((2, 5.3), (2, 2.2)),      # 6 to final (No path)
        ((4.3, 3.5), (3, 1.8)),    # Converged Yes to final
        ((5.7, 3.8), (7.5, 7.8))   # Converged No back to step 3
    ]
    
    for start, end in arrows:
        ax.annotate('', xy=end, xytext=start, 
                   arrowprops=dict(arrowstyle='->', lw=2, color='black'))
    
    # Labels
    ax.text(3.5, 2.2, 'Yes', ha='center', fontsize=10, fontweight='bold', color='green')
    ax.text(6.8, 5.8, 'No', ha='center', fontsize=10, fontweight='bold', color='red')
    
    plt.title('Bayesian Optimization Algorithm Workflow', fontsize=16, fontweight='bold', pad=20)
    plt.tight_layout()
    plt.savefig('bayesian_workflow.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'bayesian_workflow.png'

def create_results_charts():
    """Create results comparison charts (Slide 9)"""
    # Create subplots
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(14, 10))
    
    # 1. Deposition Rate Improvement
    categories = ['Baseline', 'Optimized']
    rates = [1.0, 1.4]  # 40% improvement
    colors = ['lightcoral', 'lightgreen']
    
    bars1 = ax1.bar(categories, rates, color=colors, edgecolor='black', linewidth=2)
    ax1.set_ylabel('Deposition Rate (μm/hour)', fontweight='bold')
    ax1.set_title('Deposition Rate Improvement', fontweight='bold')
    ax1.set_ylim(0, 1.6)
    
    # Add value labels on bars
    for bar, rate in zip(bars1, rates):
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height + 0.05,
                f'{rate:.1f}', ha='center', va='bottom', fontweight='bold')
    
    # Add improvement percentage
    ax1.text(1, 1.2, '+40%', ha='center', va='center', fontsize=14, 
            fontweight='bold', color='green', 
            bbox=dict(boxstyle="round,pad=0.3", facecolor='lightgreen'))
    
    # 2. Process Variability Reduction
    processes = ['Baseline', 'Optimized']
    variability = [5.0, 3.75]  # 25% reduction
    
    bars2 = ax2.bar(processes, variability, color=['lightcoral', 'lightblue'], 
                   edgecolor='black', linewidth=2)
    ax2.set_ylabel('Process Variability (%RSD)', fontweight='bold')
    ax2.set_title('Process Variability Reduction', fontweight='bold')
    ax2.set_ylim(0, 6)
    
    for bar, var in zip(bars2, variability):
        height = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width()/2., height + 0.1,
                f'{var:.2f}%', ha='center', va='bottom', fontweight='bold')
    
    ax2.text(1, 4.5, '-25%', ha='center', va='center', fontsize=14, 
            fontweight='bold', color='blue', 
            bbox=dict(boxstyle="round,pad=0.3", facecolor='lightblue'))
    
    # 3. Defect Reduction
    conditions = ['Baseline', 'Optimized']
    defects = [100, 80]  # 20% reduction
    
    bars3 = ax3.bar(conditions, defects, color=['lightcoral', 'lightyellow'], 
                   edgecolor='black', linewidth=2)
    ax3.set_ylabel('Defect Density (relative)', fontweight='bold')
    ax3.set_title('Defect Density Reduction', fontweight='bold')
    ax3.set_ylim(0, 120)
    
    for bar, defect in zip(bars3, defects):
        height = bar.get_height()
        ax3.text(bar.get_x() + bar.get_width()/2., height + 2,
                f'{defect}', ha='center', va='bottom', fontweight='bold')
    
    ax3.text(1, 90, '-20%', ha='center', va='center', fontsize=14, 
            fontweight='bold', color='orange', 
            bbox=dict(boxstyle="round,pad=0.3", facecolor='lightyellow'))
    
    # 4. Cost Reduction
    cost_categories = ['Baseline', 'Optimized']
    costs = [100, 85]  # 15% reduction
    
    bars4 = ax4.bar(cost_categories, costs, color=['lightcoral', 'lightsteelblue'], 
                   edgecolor='black', linewidth=2)
    ax4.set_ylabel('Cost per Wafer (relative)', fontweight='bold')
    ax4.set_title('Manufacturing Cost Reduction', fontweight='bold')
    ax4.set_ylim(0, 120)
    
    for bar, cost in zip(bars4, costs):
        height = bar.get_height()
        ax4.text(bar.get_x() + bar.get_width()/2., height + 2,
                f'{cost}', ha='center', va='bottom', fontweight='bold')
    
    ax4.text(1, 92, '-15%', ha='center', va='center', fontsize=14, 
            fontweight='bold', color='blue', 
            bbox=dict(boxstyle="round,pad=0.3", facecolor='lightsteelblue'))
    
    plt.suptitle('Optimization Results Summary', fontsize=16, fontweight='bold')
    plt.tight_layout()
    plt.savefig('results_charts.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'results_charts.png'

def create_doe_matrix_heatmap():
    """Create DOE matrix heatmap (Slide 6)"""
    # Sample DOE data
    np.random.seed(42)
    parameters = ['Temperature', 'Pressure', 'Power', 'Flow Rate']
    experiments = [f'Exp {i+1}' for i in range(12)]
    
    # Create DOE matrix (normalized values)
    doe_data = np.random.rand(12, 4) * 2 - 1  # Values between -1 and 1 (low, medium, high)
    
    # Create DataFrame
    df_doe = pd.DataFrame(doe_data, index=experiments, columns=parameters)
    
    # Create heatmap
    fig, ax = plt.subplots(figsize=(10, 8))
    
    # Custom colormap
    sns.heatmap(df_doe, annot=True, fmt='.2f', cmap='RdYlBu_r', center=0,
                square=True, cbar_kws={'label': 'Parameter Level\n(-1: Low, 0: Medium, +1: High)'})
    
    plt.title('Design of Experiments (DOE) Matrix', fontsize=14, fontweight='bold', pad=20)
    plt.xlabel('Process Parameters', fontweight='bold')
    plt.ylabel('Experimental Runs', fontweight='bold')
    plt.tight_layout()
    plt.savefig('doe_matrix_heatmap.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'doe_matrix_heatmap.png'

def create_validation_spider_chart():
    """Create validation metrics spider chart (Slide 10)"""
    # Metrics for baseline vs optimized
    metrics = ['Thickness\nUniformity', 'Surface\nQuality', 'Crystal\nStructure', 
              'Defect\nDensity', 'Process\nStability', 'Adhesion']
    
    # Scores (out of 10)
    baseline_scores = [6, 7, 8, 5, 6, 7]
    optimized_scores = [8, 8, 8, 8, 9, 8]
    
    # Number of variables
    N = len(metrics)
    
    # Angle for each metric
    angles = [n / float(N) * 2 * np.pi for n in range(N)]
    angles += angles[:1]  # Complete the circle
    
    # Add first value to the end to close the polygon
    baseline_scores += baseline_scores[:1]
    optimized_scores += optimized_scores[:1]
    
    # Create plot
    fig, ax = plt.subplots(figsize=(10, 10), subplot_kw=dict(projection='polar'))
    
    # Plot baseline
    ax.plot(angles, baseline_scores, 'o-', linewidth=2, label='Baseline', color='red')
    ax.fill(angles, baseline_scores, alpha=0.25, color='red')
    
    # Plot optimized
    ax.plot(angles, optimized_scores, 'o-', linewidth=2, label='Optimized', color='green')
    ax.fill(angles, optimized_scores, alpha=0.25, color='green')
    
    # Add labels
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(metrics)
    ax.set_ylim(0, 10)
    ax.set_yticks([2, 4, 6, 8, 10])
    ax.set_yticklabels(['2', '4', '6', '8', '10'])
    ax.grid(True)
    
    plt.legend(loc='upper right', bbox_to_anchor=(1.3, 1.0))
    plt.title('Film Quality Validation Metrics\n(Higher scores = Better performance)', 
             fontsize=14, fontweight='bold', pad=30)
    plt.tight_layout()
    plt.savefig('validation_spider_chart.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'validation_spider_chart.png'

def create_kla_alignment_diagram():
    """Create KLA alignment diagram (Slide 11)"""
    fig, ax = plt.subplots(1, 1, figsize=(12, 8))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis('off')
    
    # Center circle - Candidate Skills
    center_circle = Circle((5, 4), 1.5, facecolor='lightblue', edgecolor='navy', linewidth=3)
    ax.add_patch(center_circle)
    ax.text(5, 4, 'Varad Lad\nSkill Set', ha='center', va='center', 
            fontsize=12, fontweight='bold')
    
    # Surrounding KLA requirements
    kla_requirements = [
        {'pos': (2, 6.5), 'text': 'In-Situ Process\nMonitoring', 'color': 'lightgreen'},
        {'pos': (8, 6.5), 'text': 'Metrology &\nInspection', 'color': 'lightcoral'},
        {'pos': (1.5, 4), 'text': 'Sensor\nIntegration', 'color': 'lightyellow'},
        {'pos': (8.5, 4), 'text': 'Advanced\nAnalytics', 'color': 'lightpink'},
        {'pos': (2, 1.5), 'text': 'Mechanical\nDesign', 'color': 'lightsteelblue'},
        {'pos': (8, 1.5), 'text': 'Process\nOptimization', 'color': 'lightgray'}
    ]
    
    # Draw requirement boxes and connections
    for req in kla_requirements:
        box = FancyBboxPatch((req['pos'][0]-0.7, req['pos'][1]-0.5), 1.4, 1, 
                           boxstyle="round,pad=0.1", facecolor=req['color'], 
                           edgecolor='darkblue', linewidth=2)
        ax.add_patch(box)
        ax.text(req['pos'][0], req['pos'][1], req['text'], ha='center', va='center', 
                fontsize=10, fontweight='bold')
        
        # Connection lines to center
        ax.plot([req['pos'][0], 5], [req['pos'][1], 4], 'k--', alpha=0.5, linewidth=2)
    
    # Add KLA logo placeholder
    kla_box = FancyBboxPatch((4, 0.5), 2, 0.8, boxstyle="round,pad=0.1", 
                           facecolor='navy', edgecolor='black', linewidth=2)
    ax.add_patch(kla_box)
    ax.text(5, 0.9, 'KLA SensArray\nDivision', ha='center', va='center', 
            fontsize=12, fontweight='bold', color='white')
    
    plt.title('Candidate-Role Alignment Matrix', fontsize=16, fontweight='bold', pad=20)
    plt.tight_layout()
    plt.savefig('kla_alignment_diagram.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'kla_alignment_diagram.png'

def create_timeline_gantt():
    """Create project timeline Gantt chart"""
    fig, ax = plt.subplots(figsize=(12, 6))
    
    # Project phases
    phases = ['Literature Review', 'DOE Planning', 'DOE Execution', 'Model Development', 
             'Bayesian Optimization', 'Validation', 'Results Analysis']
    
    # Timeline (weeks)
    start_times = [0, 2, 4, 7, 9, 12, 14]
    durations = [3, 2, 4, 3, 4, 3, 2]
    colors = plt.cm.Set3(np.linspace(0, 1, len(phases)))
    
    # Create Gantt chart
    for i, (phase, start, duration, color) in enumerate(zip(phases, start_times, durations, colors)):
        ax.barh(i, duration, left=start, height=0.6, color=color, 
               edgecolor='black', linewidth=1)
        ax.text(start + duration/2, i, f'{duration}w', ha='center', va='center', 
               fontweight='bold', fontsize=10)
    
    ax.set_yticks(range(len(phases)))
    ax.set_yticklabels(phases)
    ax.set_xlabel('Project Timeline (Weeks)', fontweight='bold')
    ax.set_title('Deposition Rate Optimization Project Timeline', fontweight='bold', pad=20)
    ax.grid(axis='x', alpha=0.3)
    ax.set_xlim(0, 16)
    
    plt.tight_layout()
    plt.savefig('project_timeline.png', dpi=300, bbox_inches='tight')
    plt.close()
    return 'project_timeline.png'

def create_visual_presentation():
    """Create PowerPoint with all visual elements"""
    prs = Presentation()
    
    # Generate all charts
    print("Generating charts and diagrams...")
    chart_files = []
    
    chart_files.append(create_process_flow_diagram())
    chart_files.append(create_parameter_ranges_table())
    chart_files.append(create_doe_matrix_heatmap())
    chart_files.append(create_bayesian_optimization_workflow())
    chart_files.append(create_results_charts())
    chart_files.append(create_validation_spider_chart())
    chart_files.append(create_kla_alignment_diagram())
    chart_files.append(create_timeline_gantt())
    
    # Slide mapping information
    slide_mappings = [
        {
            'title': 'Process Flow Diagram',
            'description': 'Two-Phase Optimization Methodology',
            'file': 'process_flow_diagram.png',
            'main_slide': 'Slide 5: Concept and Methodology',
            'usage': 'Visual representation of the DOE + Bayesian optimization workflow'
        },
        {
            'title': 'Parameter Ranges Table',
            'description': 'Process Parameters and Optimization Results',
            'file': 'parameter_ranges_table.png',
            'main_slide': 'Slide 7: Technical Implementation & Parameter Mapping',
            'usage': 'Detailed parameter specifications and optimization outcomes'
        },
        {
            'title': 'DOE Matrix Heatmap',
            'description': 'Design of Experiments Matrix Visualization',
            'file': 'doe_matrix_heatmap.png',
            'main_slide': 'Slide 6: Decision Framework (DOE & Design Tree)',
            'usage': 'Visual representation of experimental design structure'
        },
        {
            'title': 'Bayesian Optimization Workflow',
            'description': 'Algorithm Implementation Flowchart',
            'file': 'bayesian_workflow.png',
            'main_slide': 'Slide 8: Bayesian Optimization Code and Analysis',
            'usage': 'Step-by-step algorithm workflow with decision points'
        },
        {
            'title': 'Results Summary Charts',
            'description': 'Key Performance Improvements',
            'file': 'results_charts.png',
            'main_slide': 'Slide 9: Results & Data Interpretation',
            'usage': 'Quantitative results showing 40% rate increase, 25% variability reduction'
        },
        {
            'title': 'Validation Metrics Spider Chart',
            'description': 'Film Quality Assessment Comparison',
            'file': 'validation_spider_chart.png',
            'main_slide': 'Slide 10: Metrology & Validation',
            'usage': 'Multi-dimensional quality comparison baseline vs optimized'
        },
        {
            'title': 'KLA Alignment Diagram',
            'description': 'Candidate-Role Fit Visualization',
            'file': 'kla_alignment_diagram.png',
            'main_slide': 'Slide 11: Relevance to KLA & Candidate Fit',
            'usage': 'Visual mapping of skills to KLA SensArray requirements'
        },
        {
            'title': 'Project Timeline',
            'description': 'Gantt Chart of Project Phases',
            'file': 'project_timeline.png',
            'main_slide': 'Additional/Backup slide',
            'usage': 'Project management timeline showing systematic approach'
        }
    ]
    
    # Create slides for each chart
    for i, mapping in enumerate(slide_mappings):
        # Title slide if first
        if i == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Visual Supplements\nDeposition Rate Optimization"
            subtitle = slide.placeholders[1]
            subtitle.text = "Charts, Diagrams, and Tables\nfor Technical Presentation Enhancement"
        
        # Create content slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = mapping['title']
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add image
        slide.shapes.add_picture(mapping['file'], Inches(0.5), Inches(1.2), 
                               width=Inches(9), height=Inches(6))
        
        # Add mapping information
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(9), Inches(1))
        info_frame = info_box.text_frame
        info_text = f"Usage: {mapping['usage']}\nMaps to: {mapping['main_slide']}"
        info_frame.text = info_text
        info_frame.paragraphs[0].font.size = Pt(12)
        info_frame.paragraphs[1].font.size = Pt(12)
        info_frame.paragraphs[1].font.italic = True
    
    # Create mapping summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Visual Elements Mapping Guide"
    
    content = slide.placeholders[1]
    mapping_text = ""
    for mapping in slide_mappings:
        mapping_text += f"• {mapping['title']}\n  → {mapping['main_slide']}\n  → {mapping['usage']}\n\n"
    
    content.text = mapping_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Save presentation
    prs.save('Deposition_Rate_Optimization_Visual_Supplements.pptx')
    print("Visual supplements presentation created: Deposition_Rate_Optimization_Visual_Supplements.pptx")
    
    # Create mapping reference file
    with open('visual_elements_mapping.txt', 'w') as f:
        f.write("VISUAL ELEMENTS MAPPING REFERENCE\n")
        f.write("=" * 50 + "\n\n")
        f.write("This document maps visual elements to specific slides in the main presentation.\n\n")
        
        for mapping in slide_mappings:
            f.write(f"VISUAL: {mapping['title']}\n")
            f.write(f"FILE: {mapping['file']}\n")
            f.write(f"MAPS TO: {mapping['main_slide']}\n")
            f.write(f"DESCRIPTION: {mapping['description']}\n")
            f.write(f"USAGE: {mapping['usage']}\n")
            f.write("-" * 40 + "\n\n")
    
    print("Created visual_elements_mapping.txt with detailed mapping information")
    return chart_files

if __name__ == "__main__":
    create_visual_presentation()
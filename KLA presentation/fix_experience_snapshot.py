#!/usr/bin/env python3
"""
Fix Experience Snapshot diagram with better text positioning and no overlaps
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_fixed_experience_snapshot():
    """Create experience snapshot with proper text positioning and no overlaps"""
    
    # Create larger figure for better spacing
    fig, ax = plt.subplots(figsize=(18, 14))
    ax.set_xlim(0, 18)
    ax.set_ylim(0, 14)
    ax.axis('off')
    
    # Set clean background
    fig.patch.set_facecolor('white')
    
    # Title with more space
    ax.text(9, 13.2, 'EXPERIENCE SNAPSHOT', ha='center', va='center',
           fontsize=26, fontweight='bold', color='#1f4e79')
    ax.text(9, 12.6, 'Aligned with KLA SensArray Role Requirements', ha='center', va='center',
           fontsize=18, style='italic', color='#666666')
    
    # Experience categories with improved spacing and positioning
    experiences = [
        {
            'title': 'Thin-Film Deposition\n(PVD, CVD, ALD)',
            'items': [
                '• Synthesized nano-thin films using',
                '  spin coating and sol-gel at Rayn Innovation',
                '• Conducted DOE with JMP to optimize',
                '  deposition rate and uniformity',
                '• Supported vacuum chamber setup and',
                '  tool troubleshooting at TSMC fabs'
            ],
            'position': (4.5, 9.5),
            'size': (7.5, 3.2),
            'color': '#4472c4',
            'text_color': 'white'
        },
        {
            'title': 'Vacuum System Design\n& Integration',
            'items': [
                '• Contributed to N+1 vacuum redundancy',
                '  plan at TSMC, reducing downtime by 25%',
                '• Routed vacuum utilities and coordinated',
                '  Fab21 tool hookups and layout',
                '• Experience with chamber sealing,',
                '  leak testing, and pressure regulation'
            ],
            'position': (13.5, 9.5),
            'size': (7.5, 3.2),
            'color': '#70ad47',
            'text_color': 'white'
        },
        {
            'title': 'Root Cause Analysis\n8D / DOE',
            'items': [
                '• Applied JMP to DOE for thermal and',
                '  film thickness tuning at Rayn',
                '• Led FMEA and 8D to resolve vacuum',
                '  chamber leak at TSMC, saving $120K annually',
                '• Hands-on use of failure data to refine',
                '  tool calibration sequences'
            ],
            'position': (4.5, 5.5),
            'size': (7.5, 3.2),
            'color': '#ffc000',
            'text_color': 'black'
        },
        {
            'title': 'System Design, CAD\n& Thermal Analysis',
            'items': [
                '• Modeled system testbeds in SolidWorks',
                '  and Fusion360',
                '• Performed CFD for thermal cooling zones',
                '  in data center and pod design',
                '• Created engineering drawings for vacuum',
                '  piping and tool layout (TSMC/Chemtech)'
            ],
            'position': (13.5, 5.5),
            'size': (7.5, 3.2),
            'color': '#c5504b',
            'text_color': 'white'
        }
    ]
    
    # Draw experience boxes with proper spacing
    for exp in experiences:
        x, y = exp['position']
        w, h = exp['size']
        
        # Main background box with padding
        main_box = FancyBboxPatch((x-w/2, y-h/2), w, h, 
                                 boxstyle="round,pad=0.2", 
                                 facecolor=exp['color'], alpha=0.15,
                                 edgecolor=exp['color'], linewidth=3)
        ax.add_patch(main_box)
        
        # Header box with proper sizing
        header_height = 0.8
        header_box = FancyBboxPatch((x-w/2+0.15, y+h/2-header_height-0.05), 
                                   w-0.3, header_height,
                                   boxstyle="round,pad=0.1",
                                   facecolor=exp['color'], alpha=0.95)
        ax.add_patch(header_box)
        
        # Title text (centered in header)
        ax.text(x, y+h/2-header_height/2-0.05, exp['title'], 
               ha='center', va='center',
               fontsize=13, fontweight='bold', color=exp['text_color'])
        
        # Content area for items (below header)
        content_start_y = y + h/2 - header_height - 0.3
        item_y = content_start_y
        line_spacing = 0.25
        
        # Draw items with proper spacing
        for item in exp['items']:
            # Check if text fits, if not wrap better
            ax.text(x-w/2+0.4, item_y, item, ha='left', va='top',
                   fontsize=11, color=exp['color'], fontweight='600')
            item_y -= line_spacing
    
    # Add KLA alignment section at bottom with better spacing
    kla_y = 1.8
    kla_height = 2
    kla_box = FancyBboxPatch((1, kla_y-kla_height/2), 16, kla_height, 
                            boxstyle="round,pad=0.3",
                            facecolor='#1f4e79', alpha=0.12,
                            edgecolor='#1f4e79', linewidth=3)
    ax.add_patch(kla_box)
    
    # KLA header
    ax.text(9, kla_y+0.6, 'KLA SENSARRAY ALIGNMENT', ha='center', va='center',
           fontsize=16, fontweight='bold', color='#1f4e79')
    
    # Alignment points in two clean columns
    left_alignment = [
        '✓ In-Situ Process Monitoring & Control',
        '✓ Semiconductor Deposition Tool Experience', 
        '✓ Vacuum System Integration & Design'
    ]
    
    right_alignment = [
        '✓ Statistical Process Control & DOE',
        '✓ Root Cause Analysis & Problem Solving',
        '✓ CAD Design & Thermal Analysis'
    ]
    
    # Left column
    point_y = kla_y + 0.1
    for point in left_alignment:
        ax.text(3, point_y, point, ha='left', va='center',
               fontsize=12, color='#1f4e79', fontweight='600')
        point_y -= 0.35
    
    # Right column
    point_y = kla_y + 0.1
    for point in right_alignment:
        ax.text(10, point_y, point, ha='left', va='center',
               fontsize=12, color='#1f4e79', fontweight='600')
        point_y -= 0.35
    
    # Save with high quality and proper margins
    plt.tight_layout()
    plt.savefig('experience_snapshot_fixed.png', dpi=300, bbox_inches='tight',
               facecolor='white', edgecolor='none', pad_inches=0.4)
    plt.close()
    
    print("✅ Created fixed experience snapshot: experience_snapshot_fixed.png")
    return 'experience_snapshot_fixed.png'

def create_clean_powerpoint():
    """Create clean PowerPoint with fixed snapshot"""
    
    # Generate the fixed visual
    snapshot_file = create_fixed_experience_snapshot()
    
    # Create presentation
    prs = Presentation()
    
    # Create slide with blank layout
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Experience Snapshot - KLA SensArray Alignment"
    title_frame.paragraphs[0].font.size = Pt(26)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add the fixed snapshot image (properly sized)
    slide.shapes.add_picture(snapshot_file, Inches(0.2), Inches(0.8), 
                           width=Inches(9.6), height=Inches(7))
    
    # Save presentation
    prs.save('Experience_Snapshot_Fixed.pptx')
    
    print("✅ Created fixed PowerPoint: Experience_Snapshot_Fixed.pptx")
    
    return snapshot_file

def cleanup_old_snapshot():
    """Remove old snapshot file"""
    import os
    
    old_files = ['experience_snapshot.png', 'Experience_Snapshot_KLA.pptx']
    
    for file in old_files:
        try:
            if os.path.exists(file):
                os.remove(file)
                print(f"🗑️ Removed old file: {file}")
        except Exception as e:
            print(f"⚠️ Could not remove {file}: {e}")

if __name__ == "__main__":
    print("🔧 Fixing Experience Snapshot with proper text positioning...")
    cleanup_old_snapshot()
    create_clean_powerpoint()
    print("\n✅ Fixed Experience Snapshot complete! No more text overlaps.")
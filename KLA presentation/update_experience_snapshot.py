#!/usr/bin/env python3
"""
Update Experience Snapshot with specific text changes and perfect formatting
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_updated_experience_snapshot():
    """Create experience snapshot with updated text and perfect formatting"""
    
    # Create figure with optimal size for clarity
    fig, ax = plt.subplots(figsize=(18, 14))
    ax.set_xlim(0, 18)
    ax.set_ylim(0, 14)
    ax.axis('off')
    
    # Set clean background
    fig.patch.set_facecolor('white')
    
    # Title with perfect spacing
    ax.text(9, 13.2, 'EXPERIENCE SNAPSHOT', ha='center', va='center',
           fontsize=26, fontweight='bold', color='#1f4e79')
    ax.text(9, 12.6, 'Aligned with KLA SensArray Role Requirements', ha='center', va='center',
           fontsize=18, style='italic', color='#666666')
    
    # Experience categories with updated text and perfect alignment
    experiences = [
        {
            'title': 'Thin-Film Deposition\n(PVD, CVD, ALD)',
            'items': [
                '• Synthesized nano-thin films using',
                '  spin coating and sol-gel at Rayn Innovation',
                '• Conducted DOE with JMP to optimize',
                '  deposition rate and uniformity',
                '• Supported vacuum chamber setup and',
                '  tool troubleshooting at TSMC fabs'
            ],
            'position': (4.5, 9.5),
            'size': (7.5, 3.2),
            'color': '#4472c4',
            'text_color': 'white'
        },
        {
            'title': 'Vacuum System Design\n& Integration',
            'items': [
                '• Contributed to N+1 vacuum redundancy',
                '  plan at TSMC, reducing downtime by 25%',
                '• Routed vacuum utilities and coordinated',
                '  Fab21 tool hookups and layout',
                '• Experience with chamber sealing,',
                '  leak testing, and pressure regulation'
            ],
            'position': (13.5, 9.5),
            'size': (7.5, 3.2),
            'color': '#70ad47',
            'text_color': 'white'
        },
        {
            'title': 'Root Cause Analysis\n8D / DOE',
            'items': [
                '• Applied JMP to DOE for thermal and',
                '  film thickness tuning at Rayn',
                '• Led FMEA and 8D to resolve vacuum',
                '  chamber leak at TSMC, saving $120,000',  # UPDATED TEXT
                '• Hands-on use of failure data to refine',
                '  tool calibration sequences'
            ],
            'position': (4.5, 5.5),
            'size': (7.5, 3.2),
            'color': '#ffc000',
            'text_color': 'black'
        },
        {
            'title': 'System Design, CAD\n& Thermal Analysis',
            'items': [
                '• Modeled system testbeds in SolidWorks,',
                '  ProE/Creo, Fusion360, CATIA and more',  # UPDATED TEXT
                '• Performed CFD for thermal cooling zones',
                '  in data center and pod design',
                '• Created engineering drawings for vacuum',
                '  piping and tool layout (TSMC/Chemtech)'
            ],
            'position': (13.5, 5.5),
            'size': (7.5, 3.2),
            'color': '#c5504b',
            'text_color': 'white'
        }
    ]
    
    # Draw experience boxes with perfect spacing and alignment
    for exp in experiences:
        x, y = exp['position']
        w, h = exp['size']
        
        # Main background box with proper padding
        main_box = FancyBboxPatch((x-w/2, y-h/2), w, h, 
                                 boxstyle="round,pad=0.2", 
                                 facecolor=exp['color'], alpha=0.15,
                                 edgecolor=exp['color'], linewidth=3)
        ax.add_patch(main_box)
        
        # Header box with perfect sizing
        header_height = 0.8
        header_y = y + h/2 - header_height - 0.05
        header_box = FancyBboxPatch((x-w/2+0.15, header_y), 
                                   w-0.3, header_height,
                                   boxstyle="round,pad=0.1",
                                   facecolor=exp['color'], alpha=0.95)
        ax.add_patch(header_box)
        
        # Title text (perfectly centered in header)
        title_y = header_y + header_height/2
        ax.text(x, title_y, exp['title'], 
               ha='center', va='center',
               fontsize=13, fontweight='bold', color=exp['text_color'])
        
        # Content area with optimal spacing
        content_start_y = y + h/2 - header_height - 0.45
        item_y = content_start_y
        line_spacing = 0.26
        
        # Draw items with perfect alignment
        for item in exp['items']:
            ax.text(x-w/2+0.4, item_y, item, ha='left', va='top',
                   fontsize=11, color=exp['color'], fontweight='600')
            item_y -= line_spacing
    
    # KLA alignment section with perfect formatting
    kla_y = 1.9
    kla_height = 2.2
    kla_box = FancyBboxPatch((1, kla_y-kla_height/2), 16, kla_height, 
                            boxstyle="round,pad=0.3",
                            facecolor='#1f4e79', alpha=0.12,
                            edgecolor='#1f4e79', linewidth=3)
    ax.add_patch(kla_box)
    
    # KLA header with perfect positioning
    ax.text(9, kla_y+0.7, 'KLA SENSARRAY ALIGNMENT', ha='center', va='center',
           fontsize=16, fontweight='bold', color='#1f4e79')
    
    # Alignment points in perfectly balanced columns
    left_alignment = [
        '✓ In-Situ Process Monitoring & Control',
        '✓ Semiconductor Deposition Tool Experience', 
        '✓ Vacuum System Integration & Design'
    ]
    
    right_alignment = [
        '✓ Statistical Process Control & DOE',
        '✓ Root Cause Analysis & Problem Solving',
        '✓ CAD Design & Thermal Analysis'
    ]
    
    # Left column with perfect spacing
    point_y = kla_y + 0.2
    for point in left_alignment:
        ax.text(2.8, point_y, point, ha='left', va='center',
               fontsize=12, color='#1f4e79', fontweight='600')
        point_y -= 0.35
    
    # Right column with perfect spacing
    point_y = kla_y + 0.2
    for point in right_alignment:
        ax.text(9.8, point_y, point, ha='left', va='center',
               fontsize=12, color='#1f4e79', fontweight='600')
        point_y -= 0.35
    
    # Save with optimal quality and margins
    plt.tight_layout()
    plt.savefig('experience_snapshot_updated.png', dpi=300, bbox_inches='tight',
               facecolor='white', edgecolor='none', pad_inches=0.4)
    plt.close()
    
    print("✅ Created updated experience snapshot: experience_snapshot_updated.png")
    print("📝 Applied changes:")
    print("   • '$120K annually' → '$120,000'")
    print("   • 'SolidWorks and Fusion360' → 'SolidWorks, ProE/Creo, Fusion360, CATIA and more'")
    return 'experience_snapshot_updated.png'

def create_updated_powerpoint():
    """Create PowerPoint with updated snapshot"""
    
    # Generate the updated visual
    snapshot_file = create_updated_experience_snapshot()
    
    # Create presentation
    prs = Presentation()
    
    # Create slide with blank layout
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with perfect formatting
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Experience Snapshot - KLA SensArray Alignment"
    title_frame.paragraphs[0].font.size = Pt(26)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add the updated snapshot image with perfect sizing
    slide.shapes.add_picture(snapshot_file, Inches(0.2), Inches(0.8), 
                           width=Inches(9.6), height=Inches(7))
    
    # Save presentation
    prs.save('Experience_Snapshot_Updated.pptx')
    
    print("✅ Created updated PowerPoint: Experience_Snapshot_Updated.pptx")
    
    return snapshot_file

def cleanup_previous_versions():
    """Remove previous versions"""
    import os
    
    old_files = ['experience_snapshot_fixed.png', 'Experience_Snapshot_Fixed.pptx']
    
    for file in old_files:
        try:
            if os.path.exists(file):
                os.remove(file)
                print(f"🗑️ Removed previous version: {file}")
        except Exception as e:
            print(f"⚠️ Could not remove {file}: {e}")

if __name__ == "__main__":
    print("🔄 Updating Experience Snapshot with requested text changes...")
    cleanup_previous_versions()
    create_updated_powerpoint()
    print("\n✅ Updated Experience Snapshot complete with perfect formatting!")